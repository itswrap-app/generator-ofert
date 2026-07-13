import streamlit as st
import pandas as pd
from pptx import Presentation
import gspread
from google.oauth2.service_account import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload, MediaIoBaseUpload
import io, os, subprocess, re, shutil, requests, base64, json, uuid
from pypdf import PdfWriter
from datetime import datetime
from PIL import Image
import random

# --- NOWY LINK DO CENNIKA v3 ---
LINK_DO_ARKUSZA = "https://docs.google.com/spreadsheets/d/1USF81hOinAP_vvz1QZuoNyRCT1ezJcXTDDB6RjuYtrY/edit"
PARENT_FOLDER_ID = "12HRnKn9KrZy_C1BSgv24PGD-Gl8lTRmn"

# Co ile sekund odświeżać dane z Google (cennik, szablony, rejestr, handlowcy).
# Dzięki cache aplikacja NIE odpytuje Google przy każdym kliknięciu/wpisaniu litery.
CACHE_TTL = 600

# --- BAZA OPIEKUNÓW / HANDLOWCÓW ---
# FALLBACK: używana TYLKO gdy w arkuszu nie ma zakładki "Handlowcy".
# Docelowo handlowców edytujesz w arkuszu (zakładka "Handlowcy"), kolumny:
# Handlowiec | Stanowisko | Telefon | Email | PIN | PipedriveToken | Rola
HANDLOWCY_FALLBACK = {
    "Adam Trepka": {
        "stanowisko": "CEO It`s Wrap",
        "telefon": "+48 505 008877",
        "email": "adam@itswrap.pl",
        "pin": "", "pipedrive_token": "", "rola": "admin"
    },
    "Adam Homulicki": {
        "stanowisko": "Account Manager",
        "telefon": "+48 698433834",
        "email": "adam.homulicki@itswrap.pl",
        "pin": "", "pipedrive_token": "", "rola": "handlowiec"
    },
    "Jakub Gerber": {
        "stanowisko": "Account Manager",
        "telefon": "+48 606523486",
        "email": "jakub.gerber@itswrap.pl",
        "pin": "", "pipedrive_token": "", "rola": "handlowiec"
    },
    "Daniel Heina": {
        "stanowisko": "CEO It`s Wrap",
        "telefon": "+48 609608400",
        "email": "daniel@itswrap.pl",
        "pin": "", "pipedrive_token": "", "rola": "admin"
    }
}

# --- BAZA AUT I SEGMENTÓW ---
CAR_DATABASE = {
    "Audi": {"A3": ["Hatchback", "Sedan"], "A4": ["Sedan", "Kombi"], "A6": ["Sedan", "Kombi"], "Q3": ["SUV"], "Q5": ["SUV"], "Q8": ["SUV"], "e-tron GT": ["Sedan"], "RS6": ["Kombi"]},
    "BMW": {"Seria 3": ["Sedan", "Kombi"], "Seria 4": ["Coupe", "Gran Coupe"], "Seria 5": ["Sedan", "Kombi"], "Seria 7": ["Sedan"], "X3": ["SUV"], "X5": ["SUV"], "M3": ["Sedan", "Kombi"], "M4": ["Coupe"]},
    "BYD": {"Seal": ["Sedan"], "Atto 3": ["SUV"], "Han": ["Sedan"], "Dolphin": ["Hatchback"]},
    "Ford": {"Focus": ["Hatchback", "Kombi"], "Mustang": ["Coupe", "Cabriolet"], "Mustang Mach-E": ["SUV"], "Puma": ["Crossover"]},
    "Hyundai": {"Tucson": ["SUV"], "Ioniq 5": ["Hatchback/Crossover"], "Ioniq 6": ["Sedan"], "i30": ["Hatchback", "Kombi"], "Kona": ["Crossover"]},
    "Kia": {"EV6": ["Crossover"], "Sportage": ["SUV"], "Ceed": ["Hatchback", "Kombi"], "Stinger": ["Liftback"], "Sorento": ["SUV"]},
    "Lexus": {"NX": ["SUV"], "RX": ["SUV"], "ES": ["Sedan"], "LC": ["Coupe"]},
    "Mercedes-Benz": {"Klasa C": ["Sedan", "Kombi"], "Klasa E": ["Sedan", "Kombi"], "GLC": ["SUV", "Coupe"], "GLE": ["SUV", "Coupe"], "Klasa G": ["SUV"], "AMG GT": ["Coupe"]},
    "MG": {"MG4": ["Hatchback"], "HS": ["SUV"], "ZS": ["SUV"], "Cyberster": ["Roadster"]},
    "NIO": {"ET7": ["Sedan"], "ET5": ["Sedan"], "EL7": ["SUV"]},
    "Porsche": {"911 (992)": ["Coupe", "Cabriolet"], "Taycan": ["Sedan", "Cross Turismo"], "Macan": ["SUV"], "Panamera": ["Sedan"], "Cayenne": ["SUV", "Coupe"]},
    "Renault": {"Scenic E-Tech": ["Crossover"], "Megane E-Tech": ["Hatchback"], "Austral": ["SUV"], "Clio": ["Hatchback"], "Captur": ["Crossover"]},
    "Skoda": {"Octavia": ["Liftback", "Kombi"], "Superb": ["Liftback", "Kombi"], "Kodiaq": ["SUV"], "Enyaq": ["SUV", "Coupe"]},
    "Tesla": {"Model 3": ["Sedan"], "Model Y": ["SUV"], "Model S": ["Sedan"], "Model X": ["SUV"]},
    "Toyota": {"Corolla": ["Hatchback", "Sedan", "Kombi"], "Yaris": ["Hatchback"], "RAV4": ["SUV"], "C-HR": ["Crossover"], "Camry": ["Sedan"]},
    "Volkswagen": {"Golf": ["Hatchback", "Kombi"], "Passat": ["Kombi", "Sedan"], "Arteon": ["Liftback", "Kombi"], "ID.4": ["SUV"], "Tiguan": ["SUV"]},
    "Volvo": {"XC40": ["SUV"], "XC60": ["SUV"], "XC90": ["SUV"], "V60": ["Kombi"]},
    "Inna marka...": {"Wpisz ręcznie": ["Inne"]}
}

SEGMENTY_DOMYSLNE = {
    "Audi": {"A3": "Segment C", "A4": "Segment D", "A6": "Segment D", "Q3": "Segment C", "Q5": "Segment D", "Q8": "Segment E", "e-tron GT": "Segment E", "RS6": "Segment D"},
    "BMW": {"Seria 3": "Segment D", "Seria 4": "Segment D", "Seria 5": "Segment D", "Seria 7": "Segment E", "X3": "Segment D", "X5": "Segment E", "M3": "Segment D", "M4": "Segment D"},
    "Porsche": {"911 (992)": "Segment D", "Taycan": "Segment E", "Macan": "Segment D", "Panamera": "Segment E", "Cayenne": "Segment E"},
    "Tesla": {"Model 3": "Segment D", "Model Y": "Segment D", "Model S": "Segment E", "Model X": "Segment J"}
}

# ==========================================================================
# BAZA FOLII - zaktualizowana na podstawie arkusza Google
# (źródło: 1U576lt7QG75-Jezl7E3N_uv_2znzu_ykGR8CUFnOzWg)
#
# Format wpisu koloru: "Nazwa Polska (Kod) - Nazwa Angielska"
# Dzięki temu handlowiec wybiera w swoim języku, ale w ofercie i prompcie
# do AI zachowany jest pełny opis produktu z kodem (do precyzyjnej identyfikacji).
# ==========================================================================
FOIL_GROUPS = {
    "XPEL (Folie Ochronne PPF)": {
        "Bezbarwne (Twój obecny kolor)": [
            "XPEL Ultimate Plus (Wysoki Połysk)",
            "XPEL Stealth (Mat/Satyna)"
        ],
        "XPEL Color (Zmiana Koloru PPF)": [
            "Black (Połysk)",
            "White (Połysk)",
            "Red (Połysk)",
            "Nardo Grey (Połysk)",
            "Miami Blue (Połysk)"
        ]
    },
    "3M 2080 Series": {
        "Wysoki Połysk (High Gloss)": [
            "Czarny Wysoki Połysk (HG12) - High Gloss Black",
            "Czerwony Hot Rod Wysoki Połysk (HG13) - High Gloss Hot Rod Red",
            "Palony Pomarańcz Wysoki Połysk (HG14) - High Gloss Burnt Orange",
            "Jasnożółty Wysoki Połysk (HG15) - High Gloss Bright Yellow",
            "Cytrynowy Wysoki Połysk (HG65) - High Gloss Citric Acid",
            "Zielona Modliszka Wysoki Połysk (HG336) - High Gloss Mantis Green",
            "Wojskowa Zieleń Wysoki Połysk (HG36) - High Gloss Military Green",
            "Błękit Nieba Wysoki Połysk (HG26) - High Gloss Sky Blue",
            "Burzowy Szary Wysoki Połysk (HG31) - High Gloss Storm Gray"
        ],
        "Połysk (Gloss)": [
            "Biały Połysk (G10) - Gloss White",
            "Czarny Połysk (G12) - Gloss Black",
            "Czerwony Hotrod Połysk (G13) - Gloss Hotrod Red",
            "Palony Pomarańcz Połysk (G14) - Gloss Burnt Orange",
            "Jasnożółty Połysk (G15) - Gloss Bright Yellow",
            "Zielony Jasny Połysk (G16) - Gloss Light Green",
            "Głęboki Pomarańcz Połysk (G24) - Gloss Deep Orange",
            "Słonecznikowy Połysk (G25) - Gloss Sunflower",
            "Intensywny Niebieski Połysk (G47) - Gloss Intense Blue",
            "Ognista Czerwień Połysk (G53) - Gloss Flame Red",
            "Jasnopomarańczowy Połysk (G54) - Gloss Bright Orange",
            "Wyrazisty Żółty Połysk (G55) - Gloss Lucid Yellow",
            "Błękit Nieba Połysk (G77) - Gloss Sky Blue",
            "Jasna Kość Słoniowa Połysk (G79) - Gloss Light Ivory",
            "Ciemnoczerwony Połysk (G83) - Gloss Dark Red",
            "Gorący Róż Połysk (G103) - Gloss Hot Pink",
            "Niebieski Morski Połysk (G127) - Gloss Boat Blue",
            "Antracyt Połysk (G201) - Gloss Anthracite",
            "Czerwony Metalik Połysk (G203) - Gloss Red Metallic",
            "Czarny Metalik Połysk (G212) - Gloss Black Metallic",
            "Głęboki Niebieski Metalik Połysk (G217) - Gloss Deep Blue Metallic",
            "Niebieski Metalik Połysk (G227) - Gloss Blue Metallic",
            "Niebieski Północny Połysk (G272) - Gloss Midnight Blue",
            "Czarny Żar Połysk (G282) - Gloss Ember Black",
            "Czarna Róża Połysk (GP99) - Gloss Black Rose",
            "Białe Złoto Iskrzące Połysk (GP240) - Gloss White Gold Sparkle",
            "Galaktyczna Czerń Połysk (GP292) - Gloss Galaxy Black"
        ],
        "Satyna (Satin)": [
            "Biały Satyna (S10) - Satin White",
            "Czarny Satyna (S12) - Satin Black",
            "Szary Okrętowy Satyna (S51) - Satin Battleship Gray",
            "Błękit Key West Satyna (S57) - Satin Key West",
            "Białe Aluminium Satyna (S120) - Satin White Aluminum",
            "Zielone Jabłko Satyna (S196) - Satin Apple Green",
            "Ciemnoszary Satyna (S261) - Satin Dark Gray",
            "Chmura Burzowa Satyna (S271) - Satin Thundercloud",
            "Morski Blask Satyna (S327) - Satin Ocean Shimmer",
            "Gorzki Żółty Satyna (S335) - Satin Bitter Yellow",
            "Idealny Niebieski Satyna (S347) - Satin Perfect Blue",
            "Tląca się Czerwień Satyna (S363) - Satin Smoldering Red",
            "Biała Perła Satyna (SP10) - Satin Pearl White",
            "Kameleon Wulkaniczny Satyna (SP236) - Satin Flip Volcanic Flare",
            "Mrożona Wanilia Satyna (SP240) - Satin Frozen Vanilla",
            "Czarny Złoty Pył Satyna (SP242) - Satin Gold Dust Black",
            "Wampiryczna Czerń Satyna (SP273) - Satin Vampire Red",
            "Kameleon Widmowa Perła Satyna (SP280) - Satin Flip Ghost Pearl",
            "Kameleon Psychodeliczny Satyna (SP281) - Satin Flip Psychedelic"
        ],
        "Mat (Matte)": [
            "Biały Mat (M10) - Matte White",
            "Czarny Mat (M12) - Matte Black",
            "Głęboki Mat Czarny (DM12) - Dead Matte Black",
            "Srebrny Mat (M21) - Matte Silver",
            "Wojskowa Zieleń Mat (M26) - Matte Military Green",
            "Indygo Mat (M27) - Matte Indigo",
            "Czerwony Metalik Mat (M203) - Matte Red Metallic",
            "Sosnowa Zieleń Metalik Mat (M206) - Matte Pine Green Metallic",
            "Brązowy Metalik Mat (M209) - Matte Brown Metallic",
            "Grafitowy Metalik Mat (M211) - Matte Charcoal Metallic",
            "Czarny Metalik Mat (M212) - Matte Black Metallic",
            "Niebieski Metalik Mat (M227) - Matte Blue Metallic",
            "Szare Aluminium Mat (M230) - Matte Gray Aluminum",
            "Ciemnoszary Mat (M261) - Matte Dark Gray"
        ],
        "Specjalne / Struktura": [
            "Czarny Karbon (CFS12) - Carbon Fiber Black",
            "Antracytowy Karbon (CFS201) - Carbon Fiber Anthracite",
            "Szczotkowane Aluminium (BR120) - Brushed Aluminum",
            "Szczotkowana Czerń (BR212) - Brushed Black",
            "Szczotkowana Stal (BR201) - Brushed Steel",
            "Srebrny Chrom Połysk (GC451) - Gloss Silver Chrome"
        ]
    },
    "Avery Dennison SW900": {
        "Połysk (Gloss)": [
            "Biały Połysk (SW900-101-O) - Gloss White",
            "Śnieżna Biel Połysk (SW900-110-S) - Gloss White Snow",
            "Czarny Połysk (SW900-190-O) - Gloss Black",
            "Obsydianowa Czerń Połysk (SW900-191-O) - Gloss Obsidian Black",
            "Czerwony Połysk (SW900-415-O) - Gloss Red",
            "Kardynalska Czerwień Połysk (SW900-433-O) - Gloss Cardinal Red",
            "Karminowa Czerwień Połysk (SW900-436-O) - Gloss Carmine Red",
            "Żółty Połysk (SW900-235-O) - Gloss Yellow",
            "Pomarańczowy Połysk (SW900-373-O) - Gloss Orange",
            "Szmaragdowa Zieleń Połysk (SW900-771-O) - Gloss Emerald Green",
            "Ciemnozielony Połysk (SW900-792-O) - Gloss Dark Green",
            "Niebieski Połysk (SW900-677-O) - Gloss Blue",
            "Intensywny Niebieski Połysk (SW900-667-O) - Gloss Intense Blue",
            "Szary Połysk (SW900-832-O) - Gloss Grey",
            "Kamienny Szary Połysk (SW900-821-O) - Gloss Rock Grey",
            "Ciemnoszary Połysk (SW900-865-O) - Gloss Dark Grey",
            "Imprezowy Róż Połysk (SW900-517-O) - Gloss Pool Party Pink",
            "Złoty Metalik Połysk (SW900-215-M) - Gloss Metallic Gold",
            "Czysta Czerwień Metalik Połysk (SW900-401-M) - Gloss Metallic Pure Red",
            "Wibrujący Fiolet Metalik Połysk (SW900-522-M) - Gloss Metallic Vibrant Violet",
            "Pryzmatyczny Metalik Połysk (SW900-879-M) - Gloss Metallic Pride Prismatic"
        ],
        "Satyna (Satin)": [
            "Czarny Satyna (SW900-192-M) - Satin Black",
            "Srebrny Satyna (SW900-805-M) - Satin Silver",
            "Ciemnoszary Satyna (SW900-854-M) - Satin Dark Grey",
            "Safari Złoto Satyna (SW900-260-M) - Satin Safari Gold",
            "Fioletowy Satyna (SW900-566-M) - Satin Purple",
            "Szpiegowski Szary Satyna (SW900-828-M) - Satin Spy Grey"
        ],
        "Mat (Matte)": [
            "Biały Mat (SW900-102-O) - Matte White",
            "Czarny Mat (SW900-180-O) - Matte Black",
            "Ciemnoszary Mat (SW900-856-O) - Matte Dark Grey",
            "Khaki Mat (SW900-711-O) - Matte Khaki Green",
            "Stalowy Metalik Mat (SW900-840-M) - Matte Metallic Gunmetal",
            "Grafitowy Metalik Mat (SW900-845-M) - Matte Metallic Charcoal"
        ],
        "Kameleony (ColorFlow)": [
            "Miejska Dżungla (Urban Jungle) - ColorFlow Srebrny/Zielony",
            "Ryczący Piorun (Roaring Thunder) - ColorFlow Niebieski/Czerwony",
            "Świeża Wiosna (Fresh Spring) - ColorFlow Srebrny/Złoty",
            "Wschodzące Słońce (Rising Sun) - ColorFlow Czerwony/Złoty",
            "Grań Pioruna (Lightning Ridge) - ColorFlow Zielony/Fioletowy"
        ]
    },
    "Arlon": {
        "Połysk (Gloss)": [
            "Czerwony Połysk (401) - Gloss Red",
            "Biały Połysk (402) - Gloss White",
            "Czarny Połysk (403) - Gloss Black",
            "Pomarańczowy Połysk (406) - Gloss Orange",
            "Fioletowy Metalik Połysk (407) - Gloss Purple Metallic",
            "Niebiesko-Szary Połysk (408) - Gloss Blue Grey",
            "Miętowy Połysk (409) - Gloss Mint",
            "Piaskowy Brąz Połysk (410) - Gloss Sand Brown",
            "Srebrny Metalik Połysk (411) - Gloss Silver Metallic",
            "Szary Połysk (413) - Gloss Grey",
            "Jagodowy Niebieski Połysk (414) - Gloss Berry Blue",
            "Grafitowy Metalik Połysk (416) - Gloss Charcoal Metallic",
            "Karminowa Czerwień Połysk (418) - Gloss Carmine Red",
            "Jasnożółty Połysk (421) - Gloss Bright Yellow",
            "Nocny Błękit Połysk (422) - Gloss Midnight Blue",
            "Czarna Róża Metalik Połysk (424) - Gloss Black Rose Metallic",
            "Rubinowa Czerwień Metalik Połysk (427) - Gloss Ruby Red Metallic",
            "Elektryczna Limonka Połysk (429) - Gloss Electric Lime",
            "Stealth Niebieski Metalik Połysk (433) - Gloss Stealth Blue Metallic"
        ],
        "Satyna (Satin)": [
            "Biały Satyna (452) - Satin White",
            "Czarny Satyna (453) - Satin Black",
            "Grafitowy Metalik Satyna (454) - Satin Charcoal Metallic",
            "Biała Perła Satyna (455) - Satin Pearl White",
            "Zielona Żmija Metalik Satyna (459) - Satin Viper Green Metallic",
            "Rajski Róż Metalik Satyna (460) - Satin Paradise Pink Metallic"
        ],
        "Mat (Matte)": [
            "Czarny Mat (503) - Matte Black",
            "Szary Mat (505) - Matte Grey",
            "Wojskowa Zieleń Mat (506) - Matte Military Green",
            "Aztecki Brąz Metalik Mat (508) - Matte Aztec Bronze Metallic"
        ],
        "Matowe Aluminium": [
            "Matowe Aluminium Czerwone (551) - Matte Aluminium Red",
            "Matowe Aluminium Czarna Róża (554) - Matte Aluminium Black Rose",
            "Matowe Aluminium Niebieskie (555) - Matte Aluminium Blue",
            "Matowe Aluminium Fioletowe (557) - Matte Aluminium Purple",
            "Matowe Aluminium Szare (558) - Matte Aluminium Grey"
        ]
    },
    "PWF (Platinum Wrapping Film)": {
        # Standard Line - dostępna ogólnie u dystrybutorów, klasyczne kolory
        "Standard Line - Mat (Matt)": [
            "Czarny Mat - Matt Black",
            "Biały Mat - Matt White",
            "Antracytowy Mat - Matt Anthracite",
            "Ciemnoszary Mat - Matt Dark Grey",
            "Wojskowa Zieleń Mat - Matt Military Green",
            "Niebieski Mat - Matt Blue",
            "Czerwony Mat - Matt Red"
        ],
        "Standard Line - Połysk (Gloss)": [
            "Czarny Połysk - Gloss Black",
            "Biały Połysk - Gloss White",
            "Czerwony Połysk - Gloss Red",
            "Niebieski Połysk - Gloss Blue",
            "Żółty Połysk - Gloss Yellow",
            "Pomarańczowy Połysk - Gloss Orange",
            "Zielony Połysk - Gloss Green"
        ],
        # Exclusive Line - flagowe, charakterystyczne kolory PWF (efekty metalik)
        "Exclusive Line - Mat Metalik": [
            "Mat Midnight Purple (Północny Fiolet) - Matt Midnight Purple",
            "Mat Phantom Gold (Widmowe Złoto) - Matt Phantom Gold",
            "Mat Frozen Bronze (Mrożony Brąz) - Matt Frozen Bronze",
            "Mat Smaragd (Szmaragd) - Matt Smaragd",
            "Mat Caribbean Mint (Karaibska Mięta) - Matt Caribbean Mint",
            "Mat Verdoro Green (Zieleń Verdoro) - Matt Verdoro Green",
            "Mat Tizzy Teal (Turkus Tizzy) - Matt Tizzy Teal",
            "Mat Galactic Beam (Galaktyczna Wiązka) - Matt Galactic Beam",
            "Mat Krypton Green (Zieleń Kryptonu) - Matt Krypton Green",
            "Mat Obsidian Black (Czarny Obsydian) - Matt Obsidian Black",
            "Mat Anodized Red (Anodowana Czerwień) - Matt Anodized Red",
            "Mat Anodized Blue (Anodowany Błękit) - Matt Anodized Blue",
            "Mat Ruby Red (Rubinowa Czerwień) - Matt Ruby Red"
        ],
        # Limited Edition - serie limitowane z numerem rolki (trofeum)
        "Limited Edition (numerowana seria)": [
            "Mat Royal Rose [LE] - Matt Royal Rose Limited Edition",
            "Limited Edition - inny kolor (sprawdź dostępność u PWF)"
        ],
        # Performance - nowa linia 2025: poliuretan, łączy PPF z kolorem
        "PWF Performance (PPF + kolor)": [
            "Performance - Gloss Black (PPF kolorowy)",
            "Performance - Gloss White (PPF kolorowy)",
            "Performance - Gloss Red (PPF kolorowy)",
            "Performance - Matt Black (PPF kolorowy)",
            "Performance - Matt Anthracite (PPF kolorowy)",
            "Performance - inny kolor z linii (20 kolorów)"
        ]
    },
    "Oracal 970RA": {
        "Special": [
            "Gloss Telegrey",
            "Gloss Nardo Grey Style",
            "Matte Nato Olive"
        ],
        "Metallic": [
            "Gloss Graphite Metallic",
            "Matte Anthracite Metallic",
            "Gloss Silver Grey"
        ]
    },
    "Inne (wpisz ręcznie)": {
        "Wpisz nazwę folii ręcznie": ["___CUSTOM___"]
    }
}

# ==========================================================================
# MAPOWANIE PRODUCENTA FOLII -> DOZWOLONE KATEGORIE W CENNIKU
# Logika sprzężenia panelu wizualizacji z cennikiem:
# - XPEL produkuje folie PPF (ochronne) -> w cenniku tylko kategoria "PPF"
# - 3M, Avery, Arlon, Oracal to folie wrapowe -> kategoria "Zmiana koloru"
# - "Inne" -> bez filtra, klient widzi wszystkie kategorie
#
# Specjalny przypadek: XPEL Color (Zmiana Koloru PPF) - to PPF z kolorem,
# więc dalej kategoria "PPF" (folia ochronna z efektem koloru, nie wrap winylowy).
# ==========================================================================
KATEGORIE_DLA_PRODUCENTA = {
    "XPEL (Folie Ochronne PPF)": ["PPF"],
    "3M 2080 Series": ["Zmiana koloru"],
    "Avery Dennison SW900": ["Zmiana koloru"],
    "Arlon": ["Zmiana koloru"],
    "Oracal 970RA": ["Zmiana koloru"],
    # PWF ma dwie linie produktowe:
    # - Standard/Exclusive/Limited Edition -> klasyczny wrap winylowy (Zmiana koloru)
    # - PWF Performance (od 2025) -> polyurethane łączący PPF z kolorem -> PPF
    # Domyślnie dostępne obie kategorie, klient wybiera w cenniku co pasuje.
    "PWF (Platinum Wrapping Film)": ["Zmiana koloru", "PPF"],
    "Inne (wpisz ręcznie)": None,  # None = brak filtra, wszystkie kategorie dostępne
}

# (v2.2) Migawki domyślnych baz z kodu - fallback dla pobierz_folie(),
# gdy zakładka 'Folie' w arkuszu nie istnieje lub jest pusta.
_FOIL_GROUPS_DOMYSLNE = FOIL_GROUPS
_KATEGORIE_DLA_PRODUCENTA_DOMYSLNE = KATEGORIE_DLA_PRODUCENTA

# ==========================================================================
# OPIS SCENY - CIEMNY GARAŻ Z OŚWIETLENIEM LED
# ==========================================================================
SCENE_DESCRIPTION = (
    "The car is photographed inside a completely dark modern detailing garage. "
    "Pitch black walls, black polished concrete floor with subtle reflections of the car body. "
    "The only light sources are modern cool-white LED strip lights integrated into the ceiling and along the floor edges, "
    "creating dramatic rim lighting on the car's silhouette and sharp highlights on the body panels. "
    "No other objects, no tools, no people, no text, no windows, no doors visible - only the car in pure darkness with LED lighting. "
    "Shot with professional automotive photography technique: 3/4 front angle, 35mm lens, cinematic composition, "
    "ultra sharp focus on the car, 8k resolution, photorealistic."
)


# --- FUNKCJE POMOCNICZE ---
def _czysta_nazwa_folii(folia):
    """
    Wyciąga "ludzką" nazwę folii bez kodu w nawiasie.
    Przykład: "Czarny Połysk (G12) - Gloss Black" → "Czarny Połysk"
              "XPEL Ultimate Plus (Wysoki Połysk)" → "XPEL Ultimate Plus"
    """
    return folia.split('(')[0].strip()


def install_fonts():
    font_src, font_dst = "fonts", os.path.expanduser("~/.local/share/fonts")
    if os.path.exists(font_src):
        if not os.path.exists(font_dst): os.makedirs(font_dst)
        for f in os.listdir(font_src):
            if f.lower().endswith((".ttf", ".otf")): shutil.copy(os.path.join(font_src, f), font_dst)
        subprocess.run(["fc-cache", "-f"], capture_output=True)


def _crop_to_target_ratio(img_bytes):
    """Przycina obraz do proporcji 21:18.7 (jak w oryginalnej aplikacji)."""
    img = Image.open(io.BytesIO(img_bytes))
    if img.mode != 'RGB':
        img = img.convert('RGB')
    w, h = img.size
    target_ratio = 21.0 / 18.7
    if w / h > target_ratio:
        new_w = int(h * target_ratio)
        left = (w - new_w) / 2
        img_cropped = img.crop((left, 0, left + new_w, h))
    else:
        new_h = int(w / target_ratio)
        top = (h - new_h) / 2
        img_cropped = img.crop((0, top, w, top + new_h))
    out_bytes = io.BytesIO()
    img_cropped.save(out_bytes, format='PNG')
    return out_bytes.getvalue()


def generate_ai_image(car_description, color_description, reference_images=None):
    """
    ZMIANA (v2): w razie błędu funkcja zwraca None zamiast szarego obrazka.
    Dzięki temu handlowiec NIE MOŻE nieświadomie wysłać klientowi oferty
    z szarym prostokątem zamiast auta - poprzednia wizualizacja zostaje,
    a błąd jest wyraźnie komunikowany.
    """
    api_key = st.secrets["GEMINI_API_KEY"]
    
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash-image:generateContent"
    headers = {
        "Content-Type": "application/json",
        "x-goog-api-key": api_key
    }
    
    parts = []
    
    if reference_images and len(reference_images) > 0:
        count = len(reference_images)
        if count == 1:
            ref_instruction = (
                "The attached image is a reference photograph of the EXACT car model I want to visualize. "
                "Your task: render this SAME car (identical body shape, identical headlights, identical wheels, "
                "identical proportions, identical generation/year) but place it in a new scene and change its color."
            )
        else:
            ref_instruction = (
                f"The {count} attached images are ALL reference photographs of the SAME car model taken from different angles. "
                "Use them TOGETHER to understand the exact 3D geometry, body shape, and design details of this car. "
                "Your task: render this EXACT same car model (identical generation, identical body shape, identical headlights, "
                "identical wheels, identical proportions) as ONE NEW photograph in a new scene with a new color. "
                "Do NOT try to merge or combine the images into a collage. Treat them as 3D reference of one single car."
            )
        
        full_prompt = (
            f"{ref_instruction}\n\n"
            f"NEW SCENE AND COLOR REQUIREMENTS:\n"
            f"- Car body wrap/paint color: {color_description}\n"
            f"- Scene: {SCENE_DESCRIPTION}\n\n"
            f"CRITICAL RULES:\n"
            f"- Preserve the EXACT car model from the reference(s) - same generation, same body, same face, same wheels\n"
            f"- Only change: the body color and the environment around the car\n"
            f"- Do NOT modify the car's geometry, proportions, or design details\n"
            f"- Output: ONE single photograph of the car in the described dark garage with LED lighting\n"
            f"- Framing: 3/4 front view, car centered, full body visible"
        )
        
        parts.append({"text": full_prompt})
        
        for img_bytes, mime_type in reference_images:
            base64_image = base64.b64encode(img_bytes).decode('utf-8')
            parts.append({
                "inline_data": {
                    "mime_type": mime_type,
                    "data": base64_image
                }
            })
    else:
        full_prompt = (
            f"Generate a photorealistic image of the following car:\n"
            f"{car_description}\n\n"
            f"Make sure this is the CURRENT generation of this model as sold in showrooms today - "
            f"the latest facelift or newest generation available in 2025/2026, NOT older generations. "
            f"If uncertain about the newest generation, prefer showing a generic modern body style "
            f"matching the described segment rather than an outdated specific model.\n\n"
            f"Color: the car's body is wrapped/painted in {color_description}.\n\n"
            f"Scene: {SCENE_DESCRIPTION}"
        )
        parts.append({"text": full_prompt})
    
    payload = {
        "contents": [{"parts": parts}],
        "generationConfig": {
            "responseModalities": ["IMAGE"],
            "temperature": 0.4
        }
    }
    
    try:
        response = requests.post(url, headers=headers, json=payload, timeout=120)
        
        if response.status_code != 200:
            st.error(f"API zwróciło błąd ({response.status_code}): {response.text[:800]}")
            return None
        
        data = response.json()
        
        if 'candidates' not in data or len(data['candidates']) == 0:
            st.error(f"API nie zwróciło obrazu. Pełna odpowiedź: {str(data)[:800]}")
            return None
        
        candidate = data['candidates'][0]
        
        if candidate.get('finishReason') in ('SAFETY', 'PROHIBITED_CONTENT', 'BLOCKLIST'):
            st.error(f"Generacja zablokowana przez safety filter: {candidate.get('finishReason')}")
            return None
        
        content_parts = candidate.get('content', {}).get('parts', [])
        for part in content_parts:
            inline = part.get('inlineData') or part.get('inline_data')
            if inline and inline.get('data'):
                img_data = base64.b64decode(inline['data'])
                return _crop_to_target_ratio(img_data)
        
        text_response = ""
        for part in content_parts:
            if part.get('text'):
                text_response += part['text'] + "\n"
        
        st.error(f"Model nie zwrócił obrazu. Odpowiedź modelu: {text_response[:500] or 'brak'}")
        return None
        
    except requests.exceptions.Timeout:
        st.error("Timeout - model nie odpowiedział w 120 sekund. Spróbuj ponownie.")
        return None
    except Exception as e:
        st.error(f"Błąd komunikacji z modelem: {e}")
        return None


def _wolacz_regulowy_fallback(imie):
    """
    Awaryjny generator wołacza działający OFFLINE bez AI.
    Używany gdy Gemini API nie odpowie. Nie jest idealny, ale lepszy niż mianownik.
    
    Polska gramatyka wołacza - skrót zasad:
    - Damskie na -a -> -o (Anna -> Anno, Maria -> Mario, Małgorzata -> Małgorzato)
    - Damskie na -ia -> -io (Julia -> Julio, Natalia -> Natalio)
    - Damskie na -ja -> -jo (Maja -> Majo)
    - Damskie obce/spółgłoskowe (Carmen, Beatrice, Karin) -> Pani [mianownik]
    - Męskie na -ek/-eć -> obcina się "e" (Marek->Marku, Wojciech->Wojciechu)
    - Męskie na twardą spółgłoskę -> +ie (Jan->Janie, Tomasz->Tomaszu)
    - Itd.
    """
    if not imie:
        return "Szanowny Kliencie"
    
    imie_lower = imie.lower()
    
    # ========== ŻEŃSKIE ==========
    # Imiona kończące się na -a, -ia (najczęstsze polskie damskie)
    if imie_lower.endswith('a'):
        # Wyjątki nieregularne
        wyjatki_damskie = {
            "barbara": "Barbaro", "ewa": "Ewo", "anna": "Anno",
            "maria": "Mario", "joanna": "Joanno", "magdalena": "Magdaleno",
            "agnieszka": "Agnieszko", "katarzyna": "Katarzyno",
            "małgorzata": "Małgorzato", "monika": "Moniko", "natalia": "Natalio",
            "julia": "Julio", "alicja": "Alicjo", "dorota": "Doroto",
            "renata": "Renato", "izabela": "Izabelo", "iwona": "Iwono",
            "beata": "Beato", "marta": "Marto", "paulina": "Paulino",
            "patrycja": "Patrycjo", "karolina": "Karolino", "weronika": "Weroniko",
            "aleksandra": "Aleksandro", "ola": "Olu", "kasia": "Kasiu",
            "asia": "Asiu", "basia": "Basiu", "ania": "Aniu"
        }
        if imie_lower in wyjatki_damskie:
            return f"Pani {wyjatki_damskie[imie_lower]}"
        # Reguła ogólna - damskie na -a → -o
        return f"Pani {imie[:-1]}o"
    
    # Damskie spółgłoskowe / obce (Carmen, Beatrice, Karin, Joy) - bez deklinacji
    damskie_konsonant = {"carmen", "beatrice", "karin", "joy", "ines", "doris", "ingrid"}
    if imie_lower in damskie_konsonant:
        return f"Pani {imie}"
    
    # ========== MĘSKIE ==========
    wyjatki_meskie = {
        "piotr": "Piotrze", "paweł": "Pawle", "kacper": "Kacprze",
        "marek": "Marku", "michał": "Michale", "donald": "Donaldzie",
        "konrad": "Konradzie", "dawid": "Dawidzie", "jakub": "Jakubie",
        "tomasz": "Tomaszu", "łukasz": "Łukaszu", "mateusz": "Mateuszu",
        "krzysztof": "Krzysztofie", "andrzej": "Andrzeju", "jerzy": "Jerzy",
        "wojciech": "Wojciechu", "stanisław": "Stanisławie", "władysław": "Władysławie",
        "mieczysław": "Mieczysławie", "kazimierz": "Kazimierzu",
        "robert": "Robercie", "norbert": "Norbercie", "hubert": "Hubercie",
        "albert": "Albercie", "kuba": "Kubo"  # zdrobnienie od Jakub
    }
    if imie_lower in wyjatki_meskie:
        return f"Panie {wyjatki_meskie[imie_lower]}"
    
    # Reguły końcówek
    if imie_lower.endswith('d'):
        return f"Panie {imie}zie"
    if imie_lower.endswith(('k', 'g', 'ch', 'j', 'sz', 'cz', 'rz', 'l', 'c')):
        if imie_lower.endswith('ek'):
            return f"Panie {imie[:-2]}ku"
        return f"Panie {imie}u"
    if imie_lower.endswith(('n', 'm', 'b', 'w', 'f', 's', 'z', 't', 'p')):
        return f"Panie {imie}ie"
    if imie_lower.endswith('r'):
        return f"Panie {imie}ze"
    return f"Panie {imie}"


def _zbuduj_wolacz_po_polsku(klient):
    """
    Buduje poprawną formę wołacza w języku polskim.
    
    Strategia:
    1. Próbuje przez Gemini AI - dla każdej zmienności językowej.
    2. Jeśli AI nie działa, używa fallbacku regułowego (offline słownik + reguły).
    
    AI dostaje ścisły, ograniczony prompt - musi zwrócić TYLKO gotowy wołacz,
    bez komentarzy. Walidacja po odpowiedzi sprawdza czy zaczyna się od "Pan"
    lub "Pani" - inaczej leci fallback.
    """
    imie_surowe = klient.split()[0] if klient.strip() != "" else ""
    if not imie_surowe:
        return "Szanowny Kliencie"
    
    imie = imie_surowe.title()
    
    # Próba przez AI
    try:
        api_key = st.secrets["GEMINI_API_KEY"]
    except KeyError:
        # Brak klucza - od razu fallback
        return _wolacz_regulowy_fallback(imie)
    
    prompt = (
        f"Zwróć poprawną polską formę wołacza dla imienia: {imie}\n\n"
        "Zasady:\n"
        "- Format wyjścia DOKŁADNIE: 'Panie X' (dla mężczyzny) lub 'Pani X' (dla kobiety), gdzie X to imię w wołaczu.\n"
        "- Przykłady poprawnych odmian:\n"
        "  Piotr → Panie Piotrze\n"
        "  Anna → Pani Anno\n"
        "  Maria → Pani Mario\n"
        "  Małgorzata → Pani Małgorzato\n"
        "  Julia → Pani Julio\n"
        "  Paweł → Panie Pawle\n"
        "  Marek → Panie Marku\n"
        "  Jakub → Panie Jakubie\n"
        "  Carmen → Pani Carmen (obce, bez deklinacji)\n"
        "  Quentin → Panie Quentin (obce, bez deklinacji)\n"
        "- Jeśli imię obce niedeklinowalne po polsku - zostaw w mianowniku.\n"
        "- Odpowiedz WYŁĄCZNIE samym wołaczem (np. 'Panie Piotrze'), bez cudzysłowów, bez kropki, bez komentarzy, bez wyjaśnień.\n\n"
        f"Wołacz dla '{imie}':"
    )
    
    try:
        url = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent"
        payload = {
            "contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {
                "temperature": 0.0,  # zerowa - chcemy deterministycznej, poprawnej odmiany
                "maxOutputTokens": 32,
                "thinkingConfig": {"thinkingBudget": 0}
            }
        }
        response = requests.post(
            url,
            headers={"Content-Type": "application/json", "x-goog-api-key": api_key},
            json=payload,
            timeout=10
        )
        if response.status_code != 200:
            return _wolacz_regulowy_fallback(imie)
        
        data = response.json()
        candidates = data.get('candidates', [])
        if not candidates:
            return _wolacz_regulowy_fallback(imie)
        
        parts = candidates[0].get('content', {}).get('parts', [])
        odpowiedz = ""
        for p in parts:
            if p.get('text'):
                odpowiedz += p['text']
        odpowiedz = odpowiedz.strip().strip('"\'.,!?:;')
        
        # Walidacja - musi zaczynać się od Pan/Pani i mieć sensowną długość
        if not odpowiedz or len(odpowiedz) > 50:
            return _wolacz_regulowy_fallback(imie)
        if not (odpowiedz.startswith("Panie ") or odpowiedz.startswith("Pani ")):
            return _wolacz_regulowy_fallback(imie)
        
        return odpowiedz
        
    except Exception:
        # Każdy błąd (timeout, sieć, parsing) -> regułowy fallback
        return _wolacz_regulowy_fallback(imie)


def _fallback_intro_text(wolacz, marka, czysta_folia, handlowiec_imie, handlowiec_stanowisko):
    marka_txt = marka if marka and marka != "Inna marka..." else "samochodu"
    tresc = (
        f"{wolacz},\n\n"
        f"[WSTĘP AWARYJNY - generacja AI niedostępna]\n"
        f"Dziękuję za zaufanie przy wyborze zabezpieczenia Pańskiego {marka_txt}. "
        f"Zastosowana folia {czysta_folia} zapewni skuteczną ochronę lakieru. "
        f"Zapraszam do zapoznania się ze szczegółami wyceny."
    )
    return f"{tresc}\n\nZ motoryzacyjnym pozdrowieniem,\n{handlowiec_imie}\n{handlowiec_stanowisko}"


def generate_ai_intro_text(klient, brand, model, pakiet, folia, handlowiec_imie, handlowiec_stanowisko):
    wolacz = _zbuduj_wolacz_po_polsku(klient)
    czysta_folia = _czysta_nazwa_folii(folia)
    marka_model = f"{brand} {model}".strip() if brand and brand != "Inna marka..." else "samochodu"
    
    kontekst = (
        f"Autor wiadomości: {handlowiec_imie}, {handlowiec_stanowisko}\n"
        f"Firma: ITS WRAP (premium detailing, folie ochronne PPF, oklejanie samochodów)\n"
        f"Klient: {klient if klient.strip() else 'NIEZNANE IMIĘ (zwracaj się ogólnie)'}\n"
        f"Samochód klienta: {marka_model}\n"
        f"Zamówiona usługa: {pakiet}\n"
        f"Wybrana folia: {czysta_folia}"
    )
    
    prompt = f"""Napisz krótki, spersonalizowany wstęp do oferty handlowej od specjalisty detailingu do klienta.

KONTEKST:
{kontekst}

STRUKTURA WYJŚCIA (OBOWIĄZKOWA):
Linia 1: "{wolacz},"
Linia 2: (pusta)
Linia 3+: Treść główna - DOKŁADNIE 3-4 zdania.

ZASADY TECHNICZNE:
- Odpowiedz WYŁĄCZNIE treścią wiadomości. Bez preambuły, bez komentarzy, bez markdownu, bez cudzysłowów.
- NIE dodawaj podpisu, nazwiska, stanowiska ani formuły pożegnalnej ("Pozdrawiam", "Z poważaniem" itp.) - podpis zostanie dodany automatycznie.
- NIE używaj emoji.

WYMAGANIA TREŚCIOWE (KAŻDA JEST OBOWIĄZKOWA):
1. MUSISZ wprost wspomnieć markę i model auta: {marka_model}
2. MUSISZ wprost wspomnieć nazwę folii: {czysta_folia}
3. MUSISZ podać jedną konkretną korzyść z zastosowania tej folii (ochrona lakieru, efekt wizualny, trwałość - wybierz jedną i rozwiń jednym zdaniem).
4. Ton: profesjonalny, ciepły, ludzki. Bez sprzedażowej nachalności.
5. Forma: "Pan/Pani" (ty kolumnowe). Nigdy nie przechodź na "ty".
6. ZAKAZANE frazy (są zużyte i brzmią sztucznie): "bezkompromisowe rozwiązanie", "pasja do motoryzacji", "najwyższa jakość", "na długie lata", "spokój ducha", "perfekcyjna prezencja".

PRZYKŁADY DOBREGO STYLU (do inspiracji, NIE kopiuj):
- "Dziękuję za zaufanie przy wyborze zabezpieczenia Pańskiego Porsche 911. Folia XPEL Ultimate Plus, którą zaproponowałem, skutecznie ochroni lakier przed odpryskami i zachowa fabryczny połysk przez wiele sezonów użytkowania. Zapraszam do szczegółów wyceny."
- "Miło mi przesłać wycenę zabezpieczenia BMW Serii 3. Matowe wykończenie 3M 2080 Matte Black nada autu unikalny charakter, jednocześnie chroniąc oryginalny lakier pod spodem. Wszystkie szczegóły znajduje Pan w dalszej części oferty."

Napisz teraz wstęp (zaczynając od "{wolacz},"):"""
    
    try:
        try:
            api_key = st.secrets["GEMINI_API_KEY"]
        except KeyError:
            st.error("🔑 Brak klucza GEMINI_API_KEY w secrets.toml - używam szablonu awaryjnego.")
            return _fallback_intro_text(wolacz, brand, czysta_folia, handlowiec_imie, handlowiec_stanowisko)
        
        if not api_key or len(str(api_key).strip()) < 10:
            st.error("🔑 Klucz GEMINI_API_KEY jest pusty lub nieprawidłowy - używam szablonu awaryjnego.")
            return _fallback_intro_text(wolacz, brand, czysta_folia, handlowiec_imie, handlowiec_stanowisko)
        
        url = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent"
        headers = {
            "Content-Type": "application/json",
            "x-goog-api-key": api_key
        }
        payload = {
            "contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {
                "temperature": 0.9,
                "maxOutputTokens": 1024,
                "thinkingConfig": {
                    "thinkingBudget": 0
                }
            }
        }
        
        response = requests.post(url, headers=headers, json=payload, timeout=30)
        
        if response.status_code != 200:
            st.error(
                f"❌ Gemini API zwrócił błąd {response.status_code}.\n\n"
                f"Pełna odpowiedź API:\n```\n{response.text[:1500]}\n```\n\n"
                f"Używam szablonu awaryjnego."
            )
            return _fallback_intro_text(wolacz, brand, czysta_folia, handlowiec_imie, handlowiec_stanowisko)
        
        data = response.json()
        
        if 'candidates' not in data or not data['candidates']:
            st.error(
                f"❌ Gemini nie zwrócił żadnego kandydata odpowiedzi.\n\n"
                f"Pełna odpowiedź:\n```\n{str(data)[:1500]}\n```\n\n"
                f"Używam szablonu awaryjnego."
            )
            return _fallback_intro_text(wolacz, brand, czysta_folia, handlowiec_imie, handlowiec_stanowisko)
        
        candidate = data['candidates'][0]
        finish_reason = candidate.get('finishReason')
        
        if finish_reason in ('SAFETY', 'PROHIBITED_CONTENT', 'BLOCKLIST', 'RECITATION'):
            st.error(f"❌ Gemini zablokował odpowiedź (powód: {finish_reason}) - używam szablonu awaryjnego.")
            return _fallback_intro_text(wolacz, brand, czysta_folia, handlowiec_imie, handlowiec_stanowisko)
        
        if finish_reason == 'MAX_TOKENS':
            usage = data.get('usageMetadata', {})
            thoughts = usage.get('thoughtsTokenCount', 'n/a')
            output = usage.get('candidatesTokenCount', 'n/a')
            st.error(
                f"❌ Gemini uderzył w limit tokenów. Thinking: {thoughts}, Output: {output}. "
                f"Używam szablonu awaryjnego."
            )
            return _fallback_intro_text(wolacz, brand, czysta_folia, handlowiec_imie, handlowiec_stanowisko)
        
        parts = candidate.get('content', {}).get('parts', [])
        wygenerowany_tekst = ""
        for part in parts:
            if part.get('text'):
                wygenerowany_tekst += part['text']
        
        wygenerowany_tekst = wygenerowany_tekst.strip()
        
        if not wygenerowany_tekst:
            st.error(
                f"❌ Gemini zwrócił pustą odpowiedź (finishReason: {finish_reason}).\n\n"
                f"Pełna odpowiedź kandydata:\n```\n{str(candidate)[:1500]}\n```\n\n"
                f"Używam szablonu awaryjnego."
            )
            return _fallback_intro_text(wolacz, brand, czysta_folia, handlowiec_imie, handlowiec_stanowisko)
        
        if len(wygenerowany_tekst) < 80:
            st.warning(f"⚠️ Gemini zwrócił zbyt krótki tekst ({len(wygenerowany_tekst)} znaków): '{wygenerowany_tekst}' - używam szablonu awaryjnego.")
            return _fallback_intro_text(wolacz, brand, czysta_folia, handlowiec_imie, handlowiec_stanowisko)
        
        linie = wygenerowany_tekst.split('\n')
        frazy_podpisu_start_linii = [
            'z motoryzacyjnym', 'z poważaniem', 'z wyrazami',
            'pozdrawiam', 'serdecznie pozdr', 'łączę pozdr',
            'z uszanowaniem', 'pozdrowienia', 'z najlepszymi'
        ]
        linie_wynikowe = []
        for linia in linie:
            linia_lower = linia.strip().lower()
            if any(linia_lower.startswith(fraza) for fraza in frazy_podpisu_start_linii):
                break
            linie_wynikowe.append(linia)
        
        wygenerowany_tekst = '\n'.join(linie_wynikowe).strip()
        
        if len(wygenerowany_tekst) < 80:
            st.warning(f"⚠️ Po obcięciu podpisu tekst zbyt krótki ({len(wygenerowany_tekst)} znaków) - używam szablonu awaryjnego.")
            return _fallback_intro_text(wolacz, brand, czysta_folia, handlowiec_imie, handlowiec_stanowisko)
        
        return f"{wygenerowany_tekst}\n\nZ motoryzacyjnym pozdrowieniem,\n{handlowiec_imie}\n{handlowiec_stanowisko}"
        
    except requests.exceptions.Timeout:
        st.error("⏱️ Gemini - timeout po 30s - używam szablonu awaryjnego.")
        return _fallback_intro_text(wolacz, brand, czysta_folia, handlowiec_imie, handlowiec_stanowisko)
    except Exception as e:
        import traceback
        st.error(
            f"❌ Wyjątek podczas generacji wstępu:\n```\n{type(e).__name__}: {e}\n\n{traceback.format_exc()[:1000]}\n```\n\n"
            f"Używam szablonu awaryjnego."
        )
        return _fallback_intro_text(wolacz, brand, czysta_folia, handlowiec_imie, handlowiec_stanowisko)


def download_file(service, file_id):
    request = service.files().get_media(fileId=file_id)
    fh = io.BytesIO(); downloader = MediaIoBaseDownload(fh, request)
    done = False
    while not done: _, done = downloader.next_chunk()
    fh.seek(0); return fh


def pptx_to_pdf(input_path):
    try:
        subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', os.getcwd(), input_path], check=True, capture_output=True)
        return os.path.basename(input_path).replace('.pptx', '.pdf')
    except: return None


def wybierz_strone_koncowa(pliki_na_dysku, imie_handlowca, lista_handlowcow):
    """
    Wybiera ostatnią stronę oferty spersonalizowaną pod konkretnego handlowca.
    
    Logika:
    1. Najpierw szuka pliku zawierającego imię i nazwisko handlowca w nazwie
       (np. "6_Oferta_ostatnia_Adam Trepka" dla wybranego "Adam Trepka").
    2. Jeśli brak spersonalizowanej wersji - bierze pierwszy plik z prefixem "6_"
       jako fallback ogólny.
    3. Jeśli żadnego "6_" nie ma - zwraca None (nie dodajemy strony końcowej).
    
    Porównanie nazw jest case-insensitive i ignoruje znaki interpunkcyjne typu "_"/" ",
    więc "6_Oferta_ostatnia_Adam Trepka.pptx" pasuje do handlowca "Adam Trepka".
    
    (v2: lista_handlowcow przekazywana jako parametr, bo handlowcy mogą
    pochodzić z arkusza Google, a nie ze stałej w kodzie)
    """
    def _normalizuj(tekst):
        # ujednolicamy do małych liter, zamiana _ na spację, usuwamy podwójne spacje
        return re.sub(r'\s+', ' ', tekst.lower().replace('_', ' ')).strip()
    
    imie_norm = _normalizuj(imie_handlowca)
    pliki_6 = [f for f in pliki_na_dysku if f['name'].startswith('6')]
    
    # 1. Szukamy dopasowania po imieniu handlowca
    for f in pliki_6:
        nazwa_norm = _normalizuj(f['name'])
        if imie_norm in nazwa_norm:
            return f
    
    # 2. Fallback - ogólny plik bez imienia (np. "6_Oferta_ostatnia.pptx")
    #    Wybieramy najkrótszy plik 6_ (zwykle ogólny ma najkrótszą nazwę)
    pliki_ogolne = [f for f in pliki_6 if not any(
        _normalizuj(h) in _normalizuj(f['name']) for h in lista_handlowcow
    )]
    if pliki_ogolne:
        return sorted(pliki_ogolne, key=lambda x: len(x['name']))[0]
    
    # 3. Ostateczny fallback - pierwszy lepszy plik "6_" (gdyby coś było źle nazwane)
    if pliki_6:
        return pliki_6[0]
    
    return None


# --- OBSŁUGA GOOGLE DRIVE DLA OFERT (OAuth) ---
from google.oauth2.credentials import Credentials as OAuthCredentials

def get_oauth_drive_service():
    try:
        oauth_conf = st.secrets["oauth_drive"]
    except KeyError:
        st.error(
            "🔑 Brak konfiguracji [oauth_drive] w secrets.toml. "
            "Uruchom skrypt generuj_token.py lokalnie aby wygenerować refresh token."
        )
        return None
    
    wymagane_klucze = ['client_id', 'client_secret', 'refresh_token', 'token_uri']
    brakujace = [k for k in wymagane_klucze if not oauth_conf.get(k, '').strip()]
    if brakujace:
        st.error(f"🔑 W [oauth_drive] brakuje wartości: {', '.join(brakujace)}")
        return None
    
    rt = oauth_conf['refresh_token'].strip()
    if not rt.startswith('1//') or len(rt) < 50:
        st.error(
            f"🔑 refresh_token wygląda nieprawidłowo (długość: {len(rt)} znaków, "
            f"początek: '{rt[:10]}...'). Poprawny token zaczyna się od '1//' i ma 100+ znaków. "
            f"Wygeneruj nowy używając generuj_token.py."
        )
        return None
    
    try:
        creds = OAuthCredentials(
            token=None,
            refresh_token=rt,
            token_uri=oauth_conf['token_uri'].strip(),
            client_id=oauth_conf['client_id'].strip(),
            client_secret=oauth_conf['client_secret'].strip(),
            scopes=['https://www.googleapis.com/auth/drive.file']
        )
        from google.auth.transport.requests import Request as AuthRequest
        creds.refresh(AuthRequest())
        return build('drive', 'v3', credentials=creds)
    except Exception as e:
        err_str = str(e).lower()
        if 'invalid_grant' in err_str or 'token has been expired' in err_str or 'revoked' in err_str:
            st.error(
                "🔑 **Refresh token został unieważniony przez Google.** Możliwe przyczyny:\n\n"
                "1. **OAuth consent screen jest w trybie 'Testing'** → tokeny wygasają po 7 dniach. "
                "   Wejdź w Google Cloud Console → Ekran zgody OAuth → kliknij **'Publish app'**.\n"
                "2. **Cofnąłeś dostęp aplikacji** na [myaccount.google.com/permissions](https://myaccount.google.com/permissions)\n"
                "3. **OAuth Client ID został usunięty/zmieniony** w Google Cloud Console\n"
                "4. **Hasło konta Google zostało zmienione** - inwaliduje wszystkie tokeny\n\n"
                "**Rozwiązanie:** uruchom ponownie `generuj_token.py` lokalnie, "
                "skopiuj nowe wartości do Streamlit Secrets."
            )
        else:
            st.error(f"🔑 Błąd inicjalizacji OAuth Drive: {e}")
        return None


def pobierz_lub_stworz_folder_oferty_oauth(oauth_service):
    folder_name = "Oferty ITS WRAP"
    query = f"name='{folder_name}' and mimeType='application/vnd.google-apps.folder' and trashed=false"
    results = oauth_service.files().list(
        q=query,
        fields="files(id, name)",
        spaces='drive'
    ).execute()
    items = results.get('files', [])
    
    if not items:
        file_metadata = {
            'name': folder_name,
            'mimeType': 'application/vnd.google-apps.folder'
        }
        folder = oauth_service.files().create(
            body=file_metadata,
            fields='id'
        ).execute()
        return folder.get('id')
    return items[0].get('id')


def wgraj_pdf_na_dysk(oauth_service, folder_id, file_name, file_bytes):
    try:
        file_metadata = {
            'name': file_name,
            'parents': [folder_id]
        }
        media = MediaIoBaseUpload(io.BytesIO(file_bytes), mimetype='application/pdf', resumable=True)
        file = oauth_service.files().create(
            body=file_metadata,
            media_body=media,
            fields='id, webViewLink'
        ).execute()
        
        oauth_service.permissions().create(
            fileId=file.get('id'),
            body={'type': 'anyone', 'role': 'reader'}
        ).execute()
        
        return file.get('webViewLink')
    except Exception as e:
        st.error(f"Błąd podczas zapisu pliku PDF na Google Drive: {e}")
        return None


def zapisz_do_rejestru(nr_oferty, handlowiec, klient, auto, usluga, folia, cena, pdf_link, dane_json=""):
    """
    (v2) Dodatkowy parametr dane_json - pełne parametry oferty zapisywane
    w kolumnie 'DaneOferty' (ostatnia kolumna Rejestru). Dzięki temu ofertę
    można później wczytać do edycji jednym kliknięciem.
    """
    try:
        sheet_rejestr = client.open_by_url(LINK_DO_ARKUSZA).worksheet("Rejestr")
        dzisiaj = datetime.now().strftime("%Y-%m-%d")
        nowy_wiersz = [dzisiaj, nr_oferty, handlowiec, klient, auto, usluga, folia, f"{cena} zł", "Nowa", pdf_link, dane_json]
        sheet_rejestr.append_row(nowy_wiersz)
        pobierz_rejestr.clear()  # odśwież cache rejestru, żeby nowa oferta była widoczna od razu
        return True
    except Exception as e:
        st.error(f"Nie udało się zapisać do bazy (Rejestr): {e}")
        return False


def replace_text_in_shape(shape, replacements):
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            p_text = "".join(run.text for run in p.runs)
            original_text = p_text
            for k, v in replacements.items():
                if k in p_text:
                    p_text = p_text.replace(k, str(v))
            if p_text != original_text and len(p.runs) > 0:
                p.runs[0].text = p_text
                for run in p.runs[1:]:
                    run.text = ""
                    
    if hasattr(shape, "has_table") and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                replace_text_in_shape(cell, replacements)
                
    if shape.shape_type == 6:
        for subshape in shape.shapes:
            replace_text_in_shape(subshape, replacements)


# ==========================================================================
# (v2) PIPEDRIVE - szansa sprzedaży + załączony PDF
# ==========================================================================
PIPEDRIVE_API = "https://api.pipedrive.com/v1"


def pipedrive_znajdz_lub_utworz_osobe(token, nazwa):
    """Szuka osoby po nazwie; jeśli nie istnieje - tworzy nową.
    Zwraca (person_id, czy_istniala) lub (None, False) przy błędzie."""
    try:
        r = requests.get(f"{PIPEDRIVE_API}/persons/search",
                         params={"term": nazwa, "api_token": token,
                                 "fields": "name", "exact_match": "false", "limit": 1},
                         timeout=30)
        if r.ok:
            items = (r.json().get("data") or {}).get("items") or []
            if items:
                return items[0]["item"]["id"], True
        r = requests.post(f"{PIPEDRIVE_API}/persons",
                          params={"api_token": token}, json={"name": nazwa}, timeout=30)
        if r.ok and r.json().get("success"):
            return r.json()["data"]["id"], False
        st.warning(f"Pipedrive (osoba): {r.text[:300]}")
    except Exception as e:
        st.warning(f"Pipedrive (osoba): {e}")
    return None, False


def pipedrive_utworz_szanse(token, tytul, person_id, wartosc, stage_id=None):
    """Tworzy szansę sprzedaży (deal). Zwraca deal_id lub None."""
    try:
        payload = {"title": tytul, "value": round(float(wartosc), 2), "currency": "PLN"}
        if person_id:
            payload["person_id"] = person_id
        if stage_id:
            payload["stage_id"] = stage_id
        r = requests.post(f"{PIPEDRIVE_API}/deals",
                          params={"api_token": token}, json=payload, timeout=30)
        if r.ok and r.json().get("success"):
            return r.json()["data"]["id"]
        st.warning(f"Pipedrive (szansa): {r.text[:300]}")
    except Exception as e:
        st.warning(f"Pipedrive (szansa): {e}")
    return None


@st.cache_data(ttl=CACHE_TTL, show_spinner=False)
def pipedrive_pobierz_lejki(token):
    """(v2.3) Lista lejków (pipelines) z konta handlowca: [(id, nazwa), ...]"""
    try:
        r = requests.get(f"{PIPEDRIVE_API}/pipelines",
                         params={"api_token": token}, timeout=30)
        if r.ok and r.json().get("success"):
            return [(p["id"], p["name"]) for p in (r.json().get("data") or [])]
    except Exception:
        pass
    return []


@st.cache_data(ttl=CACHE_TTL, show_spinner=False)
def pipedrive_pobierz_etapy(token, pipeline_id):
    """(v2.3) Lista etapów danego lejka: [(id, nazwa), ...] w kolejności lejka."""
    try:
        r = requests.get(f"{PIPEDRIVE_API}/stages",
                         params={"api_token": token, "pipeline_id": pipeline_id}, timeout=30)
        if r.ok and r.json().get("success"):
            return [(s["id"], s["name"]) for s in (r.json().get("data") or [])]
    except Exception:
        pass
    return []


@st.cache_data(ttl=300, show_spinner=False)
def pipedrive_szukaj_osoby(token, fraza):
    """
    (v2.6) Wyszukuje osoby w Pipedrive po nazwie i dociąga ich szczegóły
    (email, telefon, organizacja). Wynik jest cache'owany, więc wpisywanie
    klienta NIE zasypuje Pipedrive zapytaniami.
    Zwraca listę: [{"id", "name", "email", "telefon", "organizacja"}, ...]
    """
    fraza = str(fraza).strip()
    if len(fraza) < 2:
        return []
    try:
        r = requests.get(f"{PIPEDRIVE_API}/persons/search",
                         params={"term": fraza, "api_token": token,
                                 "fields": "name", "limit": 5},
                         timeout=15)
        if not (r.ok and r.json().get("success")):
            return []
        items = (r.json().get("data") or {}).get("items") or []
        wynik = []
        for it in items[:3]:  # szczegóły dociągamy dla max 3 trafień
            p = it.get("item") or {}
            pid = p.get("id")
            if not pid:
                continue
            rekord = {"id": pid, "name": p.get("name", ""),
                      "email": "", "telefon": "", "organizacja": ""}
            _org = p.get("organization")
            if isinstance(_org, dict):
                rekord["organizacja"] = _org.get("name", "") or ""
            try:
                d = requests.get(f"{PIPEDRIVE_API}/persons/{pid}",
                                 params={"api_token": token}, timeout=15)
                if d.ok and d.json().get("success"):
                    dd = d.json().get("data") or {}
                    _emaile = dd.get("email") or []
                    _tele = dd.get("phone") or []
                    if isinstance(_emaile, list):
                        rekord["email"] = next((e.get("value") for e in _emaile
                                                if isinstance(e, dict) and e.get("value")), "") or ""
                    if isinstance(_tele, list):
                        rekord["telefon"] = next((t.get("value") for t in _tele
                                                  if isinstance(t, dict) and t.get("value")), "") or ""
                    _oid = dd.get("org_id")
                    if isinstance(_oid, dict):
                        rekord["organizacja"] = _oid.get("name", "") or rekord["organizacja"]
            except Exception:
                pass
            wynik.append(rekord)
        return wynik
    except Exception:
        return []


def pipedrive_zalacz_pdf(token, deal_id, file_bytes, file_name):
    """Fizycznie załącza plik PDF do szansy sprzedaży."""
    try:
        r = requests.post(f"{PIPEDRIVE_API}/files",
                          params={"api_token": token},
                          data={"deal_id": deal_id},
                          files={"file": (file_name, file_bytes, "application/pdf")},
                          timeout=60)
        if r.ok and r.json().get("success"):
            return True
        st.warning(f"Pipedrive (plik): {r.text[:300]}")
    except Exception as e:
        st.warning(f"Pipedrive (plik): {e}")
    return False


# ==========================================================================
# (v2) POŁĄCZENIA I DANE Z GOOGLE - CACHE
# To jest lekarstwo na zawieszanie się aplikacji po każdym wpisanym polu:
# Streamlit wykonuje cały skrypt od nowa po każdej interakcji, a wcześniej
# przy każdym takim przebiegu autoryzował się w Google i pobierał cennik
# oraz listę plików. Teraz robi to raz na CACHE_TTL sekund.
# ==========================================================================
@st.cache_resource(show_spinner=False)
def get_google_clients():
    creds = Credentials.from_service_account_info(
        st.secrets["gcp_service_account"],
        scopes=["https://www.googleapis.com/auth/spreadsheets", "https://www.googleapis.com/auth/drive"]
    )
    service = build('drive', 'v3', credentials=creds)
    gs_client = gspread.authorize(creds)
    return service, gs_client


@st.cache_data(ttl=CACHE_TTL, show_spinner=False)
def pobierz_pliki_szablonow():
    service, _ = get_google_clients()
    results = service.files().list(
        q=f"'{PARENT_FOLDER_ID}' in parents and mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation' and trashed=false",
        fields="files(id, name)",
        supportsAllDrives=True,
        includeItemsFromAllDrives=True
    ).execute()
    return results.get('files', [])


@st.cache_data(ttl=CACHE_TTL, show_spinner=False)
def pobierz_cennik():
    _, gs_client = get_google_clients()
    sheet_cennik = gs_client.open_by_url(LINK_DO_ARKUSZA).worksheet("Cennik usług")
    dane = sheet_cennik.get_all_values()
    naglowki = [c.replace('\n', ' ').replace('\r', '').strip() for c in dane[0]]
    return pd.DataFrame(dane[1:], columns=naglowki)


@st.cache_data(ttl=CACHE_TTL, show_spinner=False)
def pobierz_rejestr():
    """
    (v2.1) Odporne czytanie Rejestru: zamiast get_all_records() (który wywala
    błąd, gdy wiersz ma więcej kolumn niż nagłówek - np. po dodaniu DaneOferty
    zanim dopisano nagłówek) budujemy rekordy ręcznie z get_all_values().
    """
    _, gs_client = get_google_clients()
    try:
        sheet = gs_client.open_by_url(LINK_DO_ARKUSZA).worksheet("Rejestr")
        wartosci = sheet.get_all_values()
        if len(wartosci) < 2:
            return []
        naglowki = [str(h).strip() for h in wartosci[0]]
        # Uzupełniamy brakujące/puste nagłówki technicznymi nazwami,
        # a jeśli wierszy danych jest więcej kolumn niż nagłówków -
        # dokładamy nagłówki (ostatnia "nadmiarowa" to zwykle DaneOferty).
        max_kolumn = max(len(w) for w in wartosci)
        while len(naglowki) < max_kolumn:
            naglowki.append("DaneOferty" if len(naglowki) == max_kolumn - 1 else f"Kolumna_{len(naglowki) + 1}")
        naglowki = [h if h else f"Kolumna_{i + 1}" for i, h in enumerate(naglowki)]
        rekordy = []
        for wiersz in wartosci[1:]:
            wiersz = list(wiersz) + [""] * (len(naglowki) - len(wiersz))
            rekordy.append(dict(zip(naglowki, wiersz)))
        return rekordy
    except Exception:
        return []


@st.cache_resource(show_spinner=False)
def install_fonts_once():
    """
    (v2.1) install_fonts() odpalał proces systemowy fc-cache przy KAŻDYM
    przerysowaniu strony (czyli po każdym kliknięciu w selectbox) - to było
    główne źródło "zamrażania" aplikacji. Teraz czcionki instalują się raz.
    """
    install_fonts()
    return True


@st.cache_data(ttl=CACHE_TTL, show_spinner=False)
def pobierz_handlowcow():
    """
    (v2) Handlowcy z zakładki 'Handlowcy' w arkuszu Google.
    Wymagane kolumny (1. wiersz arkusza):
    Handlowiec | Stanowisko | Telefon | Email | PIN | PipedriveToken | Rola
    - PIN: puste = brak logowania dla tej osoby (jak dotychczas)
    - PipedriveToken: osobisty token API z Pipedrive (Ustawienia -> Personal preferences -> API)
    - Rola: 'admin' widzi wszystkie oferty, 'handlowiec' tylko swoje
    Gdy zakładki nie ma - fallback do listy z kodu (pełna kompatybilność wstecz).
    """
    _, gs_client = get_google_clients()
    try:
        sheet = gs_client.open_by_url(LINK_DO_ARKUSZA).worksheet("Handlowcy")
        wynik = {}
        for row in sheet.get_all_records():
            nazwa = str(row.get("Handlowiec", "")).strip()
            if not nazwa:
                continue
            wynik[nazwa] = {
                "stanowisko": str(row.get("Stanowisko", "")).strip(),
                "telefon": str(row.get("Telefon", "")).strip(),
                "email": str(row.get("Email", "")).strip(),
                "pin": str(row.get("PIN", "")).strip(),
                "pipedrive_token": str(row.get("PipedriveToken", "")).strip(),
                "rola": str(row.get("Rola", "handlowiec")).strip().lower() or "handlowiec",
            }
        if wynik:
            return wynik, True   # dane z arkusza
        return HANDLOWCY_FALLBACK, False
    except Exception:
        return HANDLOWCY_FALLBACK, False


@st.cache_data(ttl=CACHE_TTL, show_spinner=False)
def pobierz_folie():
    """
    (v2.2) Baza folii z zakładki 'Folie' w arkuszu Google.
    Kolumny (1. wiersz): Producent | Wykończenie | Kolor | Kategorie
    - Producent:   np. "3M 2080 Series"
    - Wykończenie: np. "Połysk (Gloss)"
    - Kolor:       np. "Zielony Jasny Połysk (G16) - Gloss Light Green"
    - Kategorie:   (opcjonalna, wystarczy w pierwszym wierszu producenta)
                   kategorie cennika oddzielone przecinkiem, np. "Zmiana koloru"
                   albo "Zmiana koloru, PPF". Puste = wszystkie kategorie.
    Dodanie nowego koloru = dopisanie wiersza + przycisk 'Odśwież dane z arkusza'.
    Gdy zakładki nie ma - fallback do listy wbudowanej w kod.
    Zwraca: (foil_groups, kategorie_map, czy_z_arkusza)
    """
    _, gs_client = get_google_clients()
    try:
        sheet = gs_client.open_by_url(LINK_DO_ARKUSZA).worksheet("Folie")
        wiersze = sheet.get_all_values()
        if len(wiersze) < 2:
            return _FOIL_GROUPS_DOMYSLNE, _KATEGORIE_DLA_PRODUCENTA_DOMYSLNE, False
        grupy = {}
        kategorie_map = {}
        for w in wiersze[1:]:
            w = list(w) + [""] * (4 - len(w))
            producent = str(w[0]).strip()
            wykonczenie = str(w[1]).strip()
            kolor = str(w[2]).strip()
            kategorie_raw = str(w[3]).strip()
            if not producent or not wykonczenie or not kolor:
                continue
            grupy.setdefault(producent, {}).setdefault(wykonczenie, [])
            if kolor not in grupy[producent][wykonczenie]:
                grupy[producent][wykonczenie].append(kolor)
            if kategorie_raw and producent not in kategorie_map:
                kategorie_map[producent] = [k.strip() for k in kategorie_raw.split(",") if k.strip()]
        if not grupy:
            return _FOIL_GROUPS_DOMYSLNE, _KATEGORIE_DLA_PRODUCENTA_DOMYSLNE, False
        # Producenci bez wpisanych kategorii -> brak filtra (wszystkie kategorie)
        for producent in grupy:
            kategorie_map.setdefault(producent, None)
        # Tryb ręczny musi istnieć zawsze - dopisujemy, jeśli nie ma go w arkuszu
        if "Inne (wpisz ręcznie)" not in grupy:
            grupy["Inne (wpisz ręcznie)"] = {"Wpisz nazwę folii ręcznie": ["___CUSTOM___"]}
        kategorie_map["Inne (wpisz ręcznie)"] = None
        return grupy, kategorie_map, True
    except Exception:
        return _FOIL_GROUPS_DOMYSLNE, _KATEGORIE_DLA_PRODUCENTA_DOMYSLNE, False


def odswiez_wszystkie_dane():
    pobierz_pliki_szablonow.clear()
    pobierz_cennik.clear()
    pobierz_rejestr.clear()
    pobierz_handlowcow.clear()
    pobierz_folie.clear()


# ==========================================================================
# (v2.3) ZAPIS KONFIGURACJI Z POZIOMU APLIKACJI
# Handlowcy nie mają dostępu do arkusza - wszystkie zmiany (nowe kolory folii,
# dane handlowców) robi się w zakładce "Ustawienia" w aplikacji, a aplikacja
# sama zapisuje je do arkusza kontem technicznym.
# ==========================================================================
_NAGLOWKI_FOLIE = ["Producent", "Wykończenie", "Kolor", "Kategorie"]
_NAGLOWKI_HANDLOWCY = ["Handlowiec", "Stanowisko", "Telefon", "Email", "PIN", "PipedriveToken", "Rola"]


def _pobierz_lub_utworz_zakladke(nazwa, naglowki, rows=500, cols=10):
    _, gs_client = get_google_clients()
    ss = gs_client.open_by_url(LINK_DO_ARKUSZA)
    try:
        return ss.worksheet(nazwa), False
    except Exception:
        ws = ss.add_worksheet(title=nazwa, rows=rows, cols=cols)
        ws.update(values=[naglowki], range_name="A1")
        return ws, True


def utworz_zakladke_folie_i_importuj():
    """Tworzy zakładkę 'Folie' i wgrywa do niej CAŁĄ wbudowaną bazę kolorów."""
    ws, _ = _pobierz_lub_utworz_zakladke("Folie", _NAGLOWKI_FOLIE, rows=400, cols=6)
    wiersze = [_NAGLOWKI_FOLIE]
    for prod, cats in _FOIL_GROUPS_DOMYSLNE.items():
        if prod == "Inne (wpisz ręcznie)":
            continue
        kat = _KATEGORIE_DLA_PRODUCENTA_DOMYSLNE.get(prod)
        kat_str = ", ".join(kat) if kat else ""
        pierwszy = True
        for wyk, kolory in cats.items():
            for kolor in kolory:
                wiersze.append([prod, wyk, kolor, kat_str if pierwszy else ""])
                pierwszy = False
    ws.update(values=wiersze, range_name="A1")
    pobierz_folie.clear()


def dodaj_kolor_folii(producent, wykonczenie, kolor, kategorie_str=""):
    """Dopisuje nowy kolor do zakładki 'Folie' (tworzy ją, jeśli nie istnieje)."""
    ws, nowa = _pobierz_lub_utworz_zakladke("Folie", _NAGLOWKI_FOLIE, rows=400, cols=6)
    ws.append_row([producent.strip(), wykonczenie.strip(), kolor.strip(), kategorie_str.strip()])
    pobierz_folie.clear()
    return True


def utworz_zakladke_handlowcy_i_importuj(handlowcy_dict):
    """Tworzy zakładkę 'Handlowcy' i wgrywa aktualną listę handlowców."""
    ws, _ = _pobierz_lub_utworz_zakladke("Handlowcy", _NAGLOWKI_HANDLOWCY, rows=100, cols=10)
    wiersze = [_NAGLOWKI_HANDLOWCY]
    for nazwa, d in handlowcy_dict.items():
        wiersze.append([
            nazwa, d.get("stanowisko", ""), d.get("telefon", ""), d.get("email", ""),
            d.get("pin", ""), d.get("pipedrive_token", ""), d.get("rola", "handlowiec"),
        ])
    ws.update(values=wiersze, range_name="A1")
    pobierz_handlowcow.clear()


def zapisz_handlowca(nazwa, stanowisko, telefon, email, pin, token, rola):
    """Aktualizuje istniejącego handlowca (po nazwie) lub dodaje nowego."""
    ws, _ = _pobierz_lub_utworz_zakladke("Handlowcy", _NAGLOWKI_HANDLOWCY, rows=100, cols=10)
    wiersz_danych = [nazwa.strip(), stanowisko.strip(), telefon.strip(), email.strip(),
                     str(pin).strip(), token.strip(), (rola or "handlowiec").strip().lower()]
    wartosci = ws.get_all_values()
    for i, w in enumerate(wartosci[1:], start=2):  # start=2: numeracja wierszy arkusza
        if len(w) > 0 and str(w[0]).strip() == nazwa.strip():
            ws.update(values=[wiersz_danych], range_name=f"A{i}:G{i}")
            pobierz_handlowcow.clear()
            return "zaktualizowano"
    ws.append_row(wiersz_danych)
    pobierz_handlowcow.clear()
    return "dodano"


# --- APLIKACJA ---
st.set_page_config(
    page_title="IT'S WRAP - Generator Ofert",
    page_icon="🛡️",
    layout="wide",
    initial_sidebar_state="expanded"
)
install_fonts_once()

# Ładujemy logo
import base64 as _b64_css
def _logo_base64(sciezka):
    try:
        with open(sciezka, 'rb') as f:
            return _b64_css.b64encode(f.read()).decode('utf-8')
    except FileNotFoundError:
        return None

LOGO_PATH_CANDIDATES = ['logo.png', 'assets/logo.png', 'static/logo.png', 'images/logo.png']
LOGO_B64 = None
LOGO_EXT = None
for _cand in LOGO_PATH_CANDIDATES:
    if os.path.exists(_cand):
        LOGO_B64 = _logo_base64(_cand)
        LOGO_EXT = _cand.split('.')[-1]
        break

st.markdown("""
<link rel="preconnect" href="https://fonts.googleapis.com">
<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
<link href="https://fonts.googleapis.com/css2?family=Roboto:wght@300;400;500;700;900&display=swap" rel="stylesheet">

<style>
:root {
    --iw-blue: #007DC5;
    --iw-blue-dark: #0064A3;
    --iw-blue-darker: #005489;
    --iw-navy: #042643;
    --iw-navy-light: #003559;
    --iw-navy-lighter: #004470;
    --iw-white: #FFFFFF;
    --iw-black: #000000;
    --iw-gray-light: #F5F7FA;
    --iw-gray-border: #E1E7ED;
    --iw-gray-text: #4A5568;
}

html, body, [class*="css"], .stApp, .stMarkdown, .stText,
.stSelectbox, .stTextInput, .stNumberInput, .stButton,
.stRadio, .stCheckbox, .stFileUploader, .stDataFrame {
    font-family: 'Roboto', -apple-system, BlinkMacSystemFont, sans-serif !important;
}

h1, h2, h3, h4, h5, h6,
.stMarkdown h1, .stMarkdown h2, .stMarkdown h3 {
    font-family: 'Roboto', sans-serif !important;
    font-weight: 700 !important;
    letter-spacing: -0.01em;
    color: var(--iw-navy) !important;
}

h1, .stMarkdown h1 {
    color: var(--iw-blue) !important;
    text-transform: uppercase;
    letter-spacing: 0.02em;
    font-weight: 900 !important;
    border-bottom: 3px solid var(--iw-blue);
    padding-bottom: 12px;
    margin-bottom: 24px;
}

.stApp {
    background-color: var(--iw-white);
}

.main .block-container {
    padding-top: 2rem;
    padding-bottom: 3rem;
    max-width: 1400px;
}

section[data-testid="stSidebar"] {
    background-color: var(--iw-navy) !important;
    border-right: 4px solid var(--iw-blue);
}

section[data-testid="stSidebar"] * {
    color: var(--iw-white) !important;
}

section[data-testid="stSidebar"] h1,
section[data-testid="stSidebar"] h2,
section[data-testid="stSidebar"] h3,
section[data-testid="stSidebar"] .stMarkdown h1,
section[data-testid="stSidebar"] .stMarkdown h2 {
    color: var(--iw-white) !important;
    text-transform: uppercase;
    letter-spacing: 0.05em;
    font-weight: 700 !important;
    font-size: 1rem;
    border-bottom: 2px solid var(--iw-blue);
    padding-bottom: 6px;
    margin-top: 16px;
}

section[data-testid="stSidebar"] hr {
    border-color: rgba(0, 125, 197, 0.3) !important;
    margin: 1.5rem 0 !important;
}

section[data-testid="stSidebar"] input,
section[data-testid="stSidebar"] select,
section[data-testid="stSidebar"] textarea,
section[data-testid="stSidebar"] [data-baseweb="select"] > div {
    background-color: var(--iw-navy-light) !important;
    color: var(--iw-white) !important;
    border: 1px solid rgba(0, 125, 197, 0.4) !important;
}

section[data-testid="stSidebar"] input:focus,
section[data-testid="stSidebar"] [data-baseweb="select"]:focus-within > div {
    border-color: var(--iw-blue) !important;
    box-shadow: 0 0 0 1px var(--iw-blue) !important;
}

section[data-testid="stSidebar"] label {
    color: rgba(255, 255, 255, 0.85) !important;
    font-size: 0.85rem !important;
    font-weight: 500 !important;
}

section[data-testid="stSidebar"] [data-testid="stFileUploadDropzone"] {
    background-color: var(--iw-navy-light) !important;
    border: 2px dashed rgba(0, 125, 197, 0.5) !important;
}

section[data-testid="stSidebar"] [data-testid="stFileUploadDropzone"]:hover {
    border-color: var(--iw-blue) !important;
    background-color: var(--iw-navy-lighter) !important;
}

.stButton > button {
    background-color: var(--iw-blue) !important;
    color: var(--iw-white) !important;
    border: none !important;
    border-radius: 4px !important;
    font-weight: 700 !important;
    letter-spacing: 0.03em;
    text-transform: uppercase;
    font-size: 0.9rem !important;
    padding: 10px 24px !important;
    transition: all 0.2s ease;
    box-shadow: 0 2px 4px rgba(4, 38, 67, 0.15);
}

.stButton > button:hover {
    background-color: var(--iw-blue-dark) !important;
    box-shadow: 0 4px 12px rgba(0, 125, 197, 0.3);
    transform: translateY(-1px);
}

.stButton > button:active {
    transform: translateY(0);
    box-shadow: 0 2px 4px rgba(4, 38, 67, 0.15);
}

.stButton > button[kind="primary"] {
    background-color: var(--iw-blue) !important;
    font-weight: 900 !important;
    padding: 12px 28px !important;
}

.stButton > button[kind="primary"]:hover {
    background-color: var(--iw-navy) !important;
}

.stDownloadButton > button {
    background-color: var(--iw-navy) !important;
    color: var(--iw-white) !important;
    border: 2px solid var(--iw-blue) !important;
    font-weight: 700 !important;
    text-transform: uppercase;
}

.stDownloadButton > button:hover {
    background-color: var(--iw-blue) !important;
    border-color: var(--iw-blue) !important;
}

.main input, .main select, .main textarea,
.main [data-baseweb="select"] > div,
.main [data-baseweb="input"] {
    border: 1px solid var(--iw-gray-border) !important;
    border-radius: 4px !important;
    background-color: var(--iw-white) !important;
}

.main input:focus, .main textarea:focus,
.main [data-baseweb="select"]:focus-within > div,
.main [data-baseweb="input"]:focus-within {
    border-color: var(--iw-blue) !important;
    box-shadow: 0 0 0 2px rgba(0, 125, 197, 0.15) !important;
}

.stTabs [data-baseweb="tab-list"] {
    gap: 0;
    border-bottom: 2px solid var(--iw-gray-border);
    background: transparent;
}

.stTabs [data-baseweb="tab"] {
    background: transparent !important;
    border: none !important;
    border-radius: 0 !important;
    padding: 12px 24px !important;
    font-weight: 700 !important;
    text-transform: uppercase;
    letter-spacing: 0.03em;
    font-size: 0.9rem !important;
    color: var(--iw-gray-text) !important;
}

.stTabs [data-baseweb="tab"][aria-selected="true"] {
    color: var(--iw-blue) !important;
    border-bottom: 3px solid var(--iw-blue) !important;
    margin-bottom: -2px;
}

div[data-testid="stAlert"] {
    border-radius: 4px !important;
    border-left-width: 4px !important;
}

div[data-baseweb="notification"][kind="info"],
div[data-testid="stAlert"][class*="info"] {
    background-color: rgba(0, 125, 197, 0.08) !important;
    border-left-color: var(--iw-blue) !important;
}

div[data-testid="stAlert"][class*="success"] {
    background-color: rgba(34, 139, 89, 0.08) !important;
    border-left-color: #228B59 !important;
}

.main [data-testid="stImage"] img {
    border-radius: 4px;
    border: 1px solid var(--iw-gray-border);
    box-shadow: 0 4px 16px rgba(4, 38, 67, 0.08);
}

[data-testid="stDataFrame"], [data-testid="stDataEditor"] {
    border: 1px solid var(--iw-gray-border) !important;
    border-radius: 4px !important;
}

[data-testid="stDataFrame"] thead th,
[data-testid="stDataEditor"] thead th {
    background-color: var(--iw-navy) !important;
    color: var(--iw-white) !important;
    font-weight: 700 !important;
    text-transform: uppercase;
    font-size: 0.8rem !important;
    letter-spacing: 0.02em;
}

.iw-footer {
    margin-top: 3rem;
    padding: 1.5rem 0;
    border-top: 2px solid var(--iw-gray-border);
    text-align: center;
    color: var(--iw-gray-text);
    font-size: 0.8rem;
    text-transform: uppercase;
    letter-spacing: 0.05em;
}

.iw-footer strong {
    color: var(--iw-blue);
    font-weight: 900;
}

.iw-sidebar-logo {
    text-align: center;
    padding: 1rem 0.5rem 1.5rem 0.5rem;
    margin-bottom: 0.5rem;
    border-bottom: 2px solid var(--iw-blue);
}

.iw-sidebar-logo img {
    max-width: 180px;
    height: auto;
}

.iw-sidebar-claim {
    text-align: center;
    font-size: 0.7rem;
    letter-spacing: 0.15em;
    color: var(--iw-blue) !important;
    text-transform: uppercase;
    font-weight: 700;
    margin-top: 0.5rem;
}

.iw-main-header {
    display: flex;
    align-items: center;
    gap: 20px;
    padding: 24px 0 20px 0;
    border-bottom: 3px solid var(--iw-blue);
    margin-bottom: 32px;
}

.iw-main-header-text {
    flex: 1;
}

.iw-main-header-title {
    font-family: 'Roboto', sans-serif;
    font-weight: 900;
    font-size: 2rem;
    color: var(--iw-navy);
    text-transform: uppercase;
    letter-spacing: 0.02em;
    line-height: 1.1;
    margin: 0;
}

.iw-main-header-subtitle {
    font-family: 'Roboto', sans-serif;
    font-weight: 400;
    font-size: 0.85rem;
    color: var(--iw-blue);
    text-transform: uppercase;
    letter-spacing: 0.2em;
    margin-top: 6px;
}

#MainMenu {visibility: hidden;}
footer {visibility: hidden;}
header[data-testid="stHeader"] {
    background: transparent;
}
</style>
""", unsafe_allow_html=True)

# (v2) Wszystkie połączenia i dane pobierane przez cache - zero zbędnych
# zapytań do Google przy każdym przerysowaniu strony.
service, client = get_google_clients()
pliki_na_dysku = pobierz_pliki_szablonow()
HANDLOWCY, HANDLOWCY_Z_ARKUSZA = pobierz_handlowcow()

# (v2.2) Jeśli w arkuszu istnieje zakładka 'Folie', baza kolorów i mapowanie
# kategorii cennika pochodzą z niej. Nowy kolor dodajesz wierszem w arkuszu.
FOIL_GROUPS, KATEGORIE_DLA_PRODUCENTA, FOLIE_Z_ARKUSZA = pobierz_folie()

try:
    df_cennik = pobierz_cennik()
except Exception as e:
    st.error(f"Błąd ładowania Cennika v3. Upewnij się, że link i nazwa zakładki 'Cennik usług' są poprawne. Błąd: {e}")
    st.stop()

# (v2) Dane oferty wczytanej do edycji (z zakładki Ewidencja)
ED = st.session_state.get('edycja', {})


def _eidx(options, val, default=0):
    """(v2) Bezpieczny index do selectboxów - używany przy wczytywaniu
    zapisanej oferty do edycji. Gdy wartości nie ma na liście, wraca default."""
    options = list(options)
    try:
        return options.index(val)
    except (ValueError, TypeError):
        return default if 0 <= default < len(options) else 0


def _pole(rekord, *nazwy):
    """(v2.4) Odporne pobieranie pola z wiersza rejestru - ignoruje wielkość
    liter, spacje i znaki specjalne w nazwach nagłówków ('Nr oferty' == 'nroferty')."""
    znorm = {re.sub(r'\W', '', str(k)).lower(): v for k, v in rekord.items()}
    for n in nazwy:
        v = znorm.get(re.sub(r'\W', '', n).lower())
        if v not in (None, ""):
            return v
    return ""


def _wyciagnij_dane_oferty(rekord):
    """(v2.4) Znajduje JSON z parametrami oferty w dowolnej kolumnie wiersza -
    niezależnie od tego, jak nazwano nagłówek (DaneOferty / Dane Oferty / itd.)."""
    for v in rekord.values():
        s = str(v).strip()
        if s.startswith('{') and '"nr_o"' in s:
            return s
    return ""


def _handlowiec_wiersza(rekord):
    """(v2.5) Ustala handlowca danego wiersza rejestru:
    1) po typowych nazwach nagłówka kolumny,
    2) a gdy to zawiedzie - po zawartości: szuka w wierszu wartości,
       która jest nazwiskiem z aktualnej listy handlowców.
    Dzięki temu filtr 'moje oferty' działa niezależnie od tego,
    jak nazwano kolumnę w arkuszu."""
    v = _pole(rekord, 'Handlowiec', 'Opiekun', 'Opiekun klienta', 'Sprzedawca', 'Doradca')
    if str(v).strip():
        return str(v).strip()
    for val in rekord.values():
        s = str(val).strip()
        if s and s in HANDLOWCY:
            return s
    return ""


# --- PANEL BOCZNY ---
with st.sidebar:
    if LOGO_B64:
        st.markdown(f"""
        <div class="iw-sidebar-logo">
            <img src="data:image/{LOGO_EXT};base64,{LOGO_B64}" alt="IT'S WRAP">
            <div class="iw-sidebar-claim">Make It Change</div>
        </div>
        """, unsafe_allow_html=True)
    else:
        st.markdown("""
        <div class="iw-sidebar-logo">
            <div style="font-size: 1.5rem; font-weight: 900; letter-spacing: 0.05em; color: #FFFFFF;">
                IT'S WRAP
            </div>
            <div class="iw-sidebar-claim">Make It Change</div>
        </div>
        """, unsafe_allow_html=True)
    
    st.markdown("### Opiekun Klienta")
    wybrany_handlowiec = st.selectbox("Kto przygotowuje ofertę?", list(HANDLOWCY.keys()),
                                      index=_eidx(HANDLOWCY.keys(), ED.get('handlowiec')))
    dane_handlowca_biezacego = HANDLOWCY[wybrany_handlowiec]
    # (v2.1) Dopóki nie istnieje zakładka 'Handlowcy' w arkuszu, nie ma ról -
    # wszyscy widzą pełny rejestr (dokładnie jak w starej wersji aplikacji).
    jest_adminem = (not HANDLOWCY_Z_ARKUSZA) or dane_handlowca_biezacego.get("rola", "handlowiec") == "admin"

    # (v2.6) Token Pipedrive bieżącego handlowca - używany przy polu klienta
    # (wyszukiwanie osoby) oraz przy dodawaniu szansy sprzedaży.
    _pd_token = dane_handlowca_biezacego.get("pipedrive_token", "")
    _pd_dostepny = bool(_pd_token)

    # (v2) Prosty system logowania PIN.
    # Aktywuje się TYLKO, gdy handlowiec ma ustawiony PIN w zakładce "Handlowcy".
    # Puste pole PIN w arkuszu = brak logowania (zachowanie jak dotychczas).
    zalogowany = True
    if dane_handlowca_biezacego.get("pin"):
        pin_wpisany = st.text_input("🔑 PIN", type="password", key=f"pin_{wybrany_handlowiec}")
        zalogowany = (pin_wpisany == dane_handlowca_biezacego["pin"])
        if pin_wpisany and not zalogowany:
            st.error("Błędny PIN.")
        elif not pin_wpisany:
            st.info("Podaj PIN, aby generować oferty.")

    if st.button("🔄 Odśwież dane z arkusza"):
        odswiez_wszystkie_dane()
        st.rerun()
    
    st.markdown("---")
    st.markdown("### Studio AI")
    brand = st.selectbox("Marka", list(CAR_DATABASE.keys()),
                         index=_eidx(CAR_DATABASE.keys(), ED.get('brand')))
    
    if brand == "Inna marka...":
        custom_brand = st.text_input("Wpisz markę", value=ED.get('final_brand', ''))
        custom_model = st.text_input("Wpisz model", value=ED.get('final_model', ''))
        final_brand, final_model, body, segment_domyslny = custom_brand, custom_model, "", "Segment D"
    else:
        final_brand = brand
        final_model = st.selectbox("Model", list(CAR_DATABASE[brand].keys()),
                                   index=_eidx(CAR_DATABASE[brand].keys(), ED.get('final_model')))
        body = st.selectbox("Nadwozie", CAR_DATABASE[brand][final_model],
                            index=_eidx(CAR_DATABASE[brand][final_model], ED.get('body')))
        segment_domyslny = SEGMENTY_DOMYSLNE.get(brand, {}).get(final_model, "Segment D")
        
    _segmenty_opcje = ["Segment A", "Segment B", "Segment C", "Segment D", "Segment E", "Segment J", "Wavecamper"]
    segment_final = st.selectbox("Wybierz Segment (do wyceny)", _segmenty_opcje,
                                 index=_eidx(_segmenty_opcje, ED.get('segment', segment_domyslny),
                                             default=_segmenty_opcje.index(segment_domyslny)))
        
    _lata_opcje = [str(y) for y in range(2026, 1999, -1)]
    year = st.selectbox("Rocznik", _lata_opcje, index=_eidx(_lata_opcje, ED.get('year')))
    gen_code = st.text_input("Kod karoserii (Opcjonalnie)", value=ED.get('gen_code', ''), help="Np. G70, 992")
    
    st.markdown("---")
    st.markdown("### Folia i Kolor")
    f_brand = st.selectbox("Producent", list(FOIL_GROUPS.keys()),
                           index=_eidx(FOIL_GROUPS.keys(), ED.get('f_brand')))
    f_cat = st.selectbox("Wykończenie", list(FOIL_GROUPS[f_brand].keys()),
                         index=_eidx(FOIL_GROUPS[f_brand].keys(), ED.get('f_cat')))
    
    # Tryb "Inne" - handlowiec wpisuje nazwę folii ręcznie, kategoria w cenniku
    # odblokowuje się w pełni (bo nie wiemy do której kategorii należy "egzotyczna" folia)
    if f_brand == "Inne (wpisz ręcznie)":
        custom_foil = st.text_input(
            "Nazwa folii (wpisz)",
            value=ED.get('f_color', '') if ED.get('f_brand') == "Inne (wpisz ręcznie)" else '',
            placeholder="np. SUNTEK PPF, Hexis Skintac, własna nazwa...",
            help="Wpisz dowolną nazwę folii. Po wybraniu 'Inne' w cenniku pojawiają się wszystkie kategorie."
        )
        f_color = custom_foil if custom_foil.strip() else "Folia niestandardowa"
    else:
        f_color = st.selectbox("Kolor", FOIL_GROUPS[f_brand][f_cat],
                               index=_eidx(FOIL_GROUPS[f_brand][f_cat], ED.get('f_color')))

    paint_color = ""
    if "Bezbarwne" in f_cat:
        paint_color = st.text_input("🚘 Podaj obecny kolor lakieru auta",
                                    value=ED.get('paint_color', "Czarny metallic"))

    st.markdown("---")
    st.markdown("### Wizualizacja AI")
    
    st.caption(
        "💡 **Rekomendacja:** wgraj 1-3 zdjęcia auta z konfiguratora producenta lub press-kitu "
        "(różne ujęcia tego samego modelu). Model AI zachowa bryłę i zmieni tylko kolor na ciemny garaż."
    )
    
    uploaded_files = st.file_uploader(
        "Zdjęcia referencyjne (opcjonalnie, ale mocno zalecane dla nowych modeli)", 
        type=['png', 'jpg', 'jpeg'], 
        accept_multiple_files=True
    )
    
    url_zdjecia = st.text_input(
        "...lub wklej URL zdjęcia (np. z konfiguratora)",
        help="Bezpośredni link do obrazka .jpg/.png ze strony producenta"
    )

    if st.button("🪄 GENERUJ WIZUALIZACJĘ AI", type="primary", disabled=not zalogowany):
        extra_code = f" ({gen_code})" if gen_code else ""
        car_description = (
            f"{year} {final_brand} {final_model}{extra_code}, body type: {body}. "
            f"This is a {segment_final.replace('Segment ', 'segment ')} vehicle."
        )
        
        # Opis koloru dla AI - bierzemy "czystą" nazwę bez kodu w nawiasie
        # (kod jak G12 jest niepotrzebny dla modelu obrazowego)
        czysta_nazwa_koloru = _czysta_nazwa_folii(f_color)
        
        if "Bezbarwne" in f_cat:
            if "Stealth" in f_color:
                color_description = (
                    f"the original {paint_color} factory paint covered with matte/satin transparent PPF film "
                    f"(XPEL Stealth) giving the paint a deep matte finish while keeping the original color visible"
                )
            else:
                color_description = (
                    f"the original {paint_color} factory paint covered with high-gloss transparent PPF film "
                    f"(XPEL Ultimate Plus) giving the paint extra depth and shine"
                )
        elif f_brand == "Inne (wpisz ręcznie)":
            # Tryb własnej folii - nie wiemy jaki producent, opis bardziej ogólny
            color_description = f"{czysta_nazwa_koloru} vinyl wrap"
        else:
            color_description = f"{czysta_nazwa_koloru} vinyl wrap by {f_brand}"
        
        reference_images = []
        
        if uploaded_files:
            for uf in uploaded_files:
                img_bytes = uf.read()
                mime = uf.type if uf.type else "image/jpeg"
                reference_images.append((img_bytes, mime))
        
        if url_zdjecia and url_zdjecia.strip():
            try:
                resp = requests.get(url_zdjecia.strip(), timeout=15, headers={
                    'User-Agent': 'Mozilla/5.0'
                })
                if resp.status_code == 200:
                    mime = resp.headers.get('content-type', 'image/jpeg').split(';')[0].strip()
                    if 'image' in mime:
                        reference_images.append((resp.content, mime))
                    else:
                        st.warning(f"URL nie wskazuje na obrazek (typ: {mime}).")
                else:
                    st.warning(f"Nie mogę pobrać obrazka z URL (kod: {resp.status_code}).")
            except Exception as e:
                st.warning(f"Nie udało się pobrać zdjęcia z URL: {e}")
        
        if reference_images:
            st.info(f"🎯 Tryb edycji: używam {len(reference_images)} zdjęcia/zdjęć referencyjnych auta.")
        else:
            st.info("🎨 Tryb generacji: tworzę auto od zera na podstawie opisu (jakość może być niższa dla bardzo nowych modeli - rozważ dodanie zdjęcia referencyjnego).")
        
        with st.spinner("Gemini 2.5 Flash Image renderuje Twoje auto w ciemnym garażu LED..."):
            img_data = generate_ai_image(
                car_description=car_description,
                color_description=color_description,
                reference_images=reference_images if reference_images else None
            )
            if img_data:
                st.session_state['ai_img'] = img_data
            else:
                # (v2) Nie nadpisujemy poprzedniej wizualizacji szarym obrazkiem.
                st.warning("Nie wygenerowano nowej wizualizacji. Poprzednia (jeśli była) pozostaje bez zmian.")
                
    st.markdown("---")
    st.markdown("### Dodatki do oferty")
    dodatki_dostepne = [f for f in pliki_na_dysku if f['name'].startswith(('4','5'))]
    wybrane_dodatki = [d for d in sorted(dodatki_dostepne, key=lambda x: x['name'])
                       if st.checkbox(d['name'], value=(d['name'] in ED.get('dodatki', [])), key=f"dod_{d['id']}")]


# ZAKŁADKI W GŁÓWNYM PANELU
tab_kreator, tab_rejestr, tab_ustawienia = st.tabs(["⚙️ Kreator Ofert", "📋 Ewidencja (Rejestr)", "🛠 Ustawienia"])

with tab_kreator:
    st.markdown("""
    <div class="iw-main-header">
        <div class="iw-main-header-text">
            <div class="iw-main-header-title">Generator Ofert</div>
            <div class="iw-main-header-subtitle">Professional Car Wrapping &amp; PPF · IT'S WRAP</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # (v2) Baner trybu edycji
    if ED:
        col_ed1, col_ed2 = st.columns([3, 1])
        with col_ed1:
            st.info(f"✏️ **Tryb edycji oferty {ED.get('nr_o', '')}** ({ED.get('klient', '')}) - "
                    f"zmień co trzeba i wygeneruj PDF ponownie. Nowa wersja zostanie dopisana do rejestru.")
        with col_ed2:
            if st.button("❌ Zakończ edycję"):
                st.session_state.pop('edycja', None)
                st.session_state.pop('ai_img', None)
                st.session_state['wstep_editor'] = ''
                st.rerun()

    col1, col2 = st.columns(2)

    with col1:
        klient = st.text_input("Imię i Nazwisko Klienta", value=ED.get('klient', ''))

        # (v2) Weryfikacja klienta w bazie - działa na danych z cache,
        # więc NIE powoduje żadnych dodatkowych zapytań ani zawieszeń.
        if klient.strip() and not ED:
            _rej = pobierz_rejestr()
            _trafienia = [r for r in _rej
                          if str(_pole(r, 'Klient')).strip().lower() == klient.strip().lower()]
            if _trafienia:
                _ost = _trafienia[-1]
                st.warning(
                    f"⚠️ Klient **{klient}** jest już w rejestrze "
                    f"(ostatnia oferta: {_pole(_ost, 'Nr oferty', 'Numer oferty', 'Nr') or '?'} "
                    f"z {_pole(_ost, 'Data') or '?'}, "
                    f"handlowiec: {_handlowiec_wiersza(_ost) or '?'}). "
                    f"Jeśli chcesz kontynuować tamtą ofertę, wczytaj ją w zakładce Ewidencja."
                )

        # (v2.6) Wyszukiwanie klienta w PIPEDRIVE - żeby nie dublować osób.
        # Jeśli osoba istnieje, jej dane (email/telefon/firma) zapisują się
        # w danych oferty, a szansa sprzedaży powiąże się z nią po ID.
        _pd_osoba = None
        if klient.strip() and _pd_dostepny:
            _pd_znalezieni = pipedrive_szukaj_osoby(_pd_token, klient.strip())
            if _pd_znalezieni:
                _pd_osoba = _pd_znalezieni[0]
                _szczegoly = " · ".join(x for x in [_pd_osoba.get("email"),
                                                    _pd_osoba.get("telefon"),
                                                    _pd_osoba.get("organizacja")] if x)
                st.success(
                    f"🟢 Pipedrive: znaleziono osobę **{_pd_osoba['name']}**"
                    + (f" ({_szczegoly})" if _szczegoly else "")
                    + " - szansa sprzedaży zostanie powiązana z nią, bez tworzenia duplikatu."
                    + (f" Podobnych osób: {len(_pd_znalezieni)}." if len(_pd_znalezieni) > 1 else "")
                )
            else:
                st.caption("⚪ Pipedrive: brak takiej osoby - zostanie utworzona przy dodawaniu szansy sprzedaży.")

        nr_o = st.text_input("Numer oferty",
                             value=ED.get('nr_o', f"IW/{datetime.now().strftime('%Y/%m/%d')}/01"))
        
        # Filtrowanie kategorii w cenniku na podstawie wybranego producenta folii
        # (XPEL -> tylko PPF, 3M/Avery/Arlon/Oracal -> tylko Zmiana koloru, Inne -> wszystkie)
        kategorie_wszystkie = [k for k in df_cennik['Kategoria'].unique() if str(k).strip() != ""]
        dozwolone_kategorie = KATEGORIE_DLA_PRODUCENTA.get(f_brand)
        
        if dozwolone_kategorie is None:
            # Tryb "Inne" - pokazujemy wszystko
            kategorie = kategorie_wszystkie
        else:
            # Filtrujemy kategorie - bierzemy tylko te które są w mapowaniu I istnieją w cenniku
            # case-insensitive, żeby działało jeśli ktoś w arkuszu wpisze "ppf" zamiast "PPF"
            dozwolone_lower = [d.lower().strip() for d in dozwolone_kategorie]
            kategorie = [k for k in kategorie_wszystkie if k.lower().strip() in dozwolone_lower]
            
            # Bezpiecznik: jeśli filtrowanie odrzuciło wszystko (np. w cenniku nie ma kolumny PPF),
            # pokazujemy wszystkie kategorie z ostrzeżeniem, żeby aplikacja się nie zablokowała
            if not kategorie:
                st.warning(
                    f"⚠️ Dla producenta '{f_brand}' oczekiwane kategorie to {dozwolone_kategorie}, "
                    f"ale w cenniku ich nie ma. Pokazuję wszystkie kategorie."
                )
                kategorie = kategorie_wszystkie
        
        kategoria = st.selectbox("Kategoria", kategorie, index=_eidx(kategorie, ED.get('kategoria')))
        uslugi_kat = [u for u in df_cennik[df_cennik['Kategoria'] == kategoria]['Usługa'].unique() if str(u).strip() != ""]
        pakiet = st.selectbox("Usługa", uslugi_kat, index=_eidx(uslugi_kat, ED.get('pakiet')))
        
        try:
            wiersz_ceny = df_cennik[(df_cennik['Kategoria'] == kategoria) & (df_cennik['Usługa'] == pakiet) & (df_cennik['Segment'] == segment_final)]
            if not wiersz_ceny.empty:
                cena_str = str(wiersz_ceny['Cena sprzedaży netto PLN'].values[0])
                cena_str = cena_str.replace(' ', '').replace('\xa0', '')
                if ',' in cena_str: cena_str = cena_str.replace('.', '').replace(',', '.')
                cena_domyslna = float(re.sub(r'[^\d.]', '', cena_str))
            else:
                cena_domyslna = 0.0
        except:
            cena_domyslna = 0.0
        
        # (v2) Cicha cena 0.0 bywa przeoczana - dajemy delikatne ostrzeżenie
        if cena_domyslna == 0.0:
            st.caption("⚠️ Nie znaleziono ceny w cenniku dla tej kombinacji usługi i segmentu - wpisz cenę ręcznie poniżej.")

        st.markdown("---")
        st.write("💰 **Kalkulacja cenowa**")
        
        cena_manual = st.number_input("Cena bazowa NETTO (PLN) - możesz edytować",
                                      value=float(ED.get('cena_manual', cena_domyslna)), step=100.0)
        rabat = st.number_input("Rabat dla klienta (PLN)",
                                value=float(ED.get('rabat', 0.0)), step=100.0)
        cena_koncowa = cena_manual - rabat
        
        st.info(f"**Cena do zapłaty netto (na ofercie): {cena_koncowa:,.2f} zł**".replace(',', 'X').replace('.', ',').replace('X', ' '))

    with col2:
        if 'ai_img' in st.session_state:
            st.image(st.session_state['ai_img'], use_container_width=True)
        else:
            st.info("Skonfiguruj auto w panelu bocznym i wygeneruj zdjęcie, aby zobaczyć podgląd.")

    # ==========================================================
    # (v2.1) TEKST WSTĘPU - podgląd i ręczna edycja przed PDF-em
    # Działa jak wizualizacja: klikasz "Generuj tekst", AI proponuje
    # treść, a Ty możesz ją dowolnie poprawić w polu poniżej.
    # Do PDF-a trafia DOKŁADNIE to, co widzisz w tym polu.
    # Jeśli pole zostawisz puste - tekst wygeneruje się automatycznie
    # w momencie tworzenia PDF (zachowanie jak dotychczas).
    # ==========================================================
    st.markdown("---")
    st.markdown("### ✍️ Tekst wstępu (strona 2 oferty)")
    # (v2.4) Wczytanie wstępu z edytowanej oferty - flaga ustawiana przy kliknięciu
    # "Edytuj" w Ewidencji, konsumowana tutaj PRZED utworzeniem pola tekstowego
    # (Streamlit nie pozwala zmieniać stanu widgetu po jego utworzeniu w tym samym przebiegu).
    if 'zaladuj_wstep' in st.session_state:
        st.session_state['wstep_editor'] = st.session_state.pop('zaladuj_wstep')
    st.session_state.setdefault('wstep_editor', '')

    col_w1, col_w2 = st.columns([1, 3])
    with col_w1:
        if st.button("🪄 GENERUJ TEKST WSTĘPU", disabled=not zalogowany):
            _autor_imie = "Adam Trepka"
            _autor_stan = HANDLOWCY.get(_autor_imie, {}).get("stanowisko", "CEO It`s Wrap")
            _folia_do_wstepu = f"{f_color} (na lakier: {paint_color})" if "Bezbarwne" in f_cat else f_color
            with st.spinner("AI pisze wstęp..."):
                st.session_state['wstep_editor'] = generate_ai_intro_text(
                    klient, final_brand, final_model, pakiet, _folia_do_wstepu,
                    _autor_imie, _autor_stan
                )
            st.rerun()
        if st.session_state.get('wstep_editor', '').strip():
            if st.button("🗑️ Wyczyść (auto przy PDF)"):
                st.session_state['wstep_editor'] = ''
                st.rerun()
    with col_w2:
        wstep_recznie = st.text_area(
            "Treść wstępu - możesz edytować przed wygenerowaniem PDF",
            key='wstep_editor',
            height=220,
            placeholder=("Pole puste = wstęp wygeneruje się automatycznie podczas tworzenia PDF.\n"
                         "Kliknij 'GENERUJ TEKST WSTĘPU', aby zobaczyć i poprawić treść przed ofertą.")
        )

    if st.button("🔥 GENERUJ PEŁNĄ OFERTĘ PDF", disabled=not zalogowany):
        if 'ai_img' not in st.session_state:
            st.error("Wizualizacja auta jest wymagana. Użyj przycisku w panelu bocznym!")
        else:
            with st.spinner("Składam profesjonalny PDF i wgrywam na Dysk Google..."):
                writer = PdfWriter()
                # W ofercie pokazujemy pełną nazwę koloru z kodem - klient widzi co konkretnie kupuje
                final_foil_text = f"{f_color} (na lakier: {paint_color})" if "Bezbarwne" in f_cat else f_color
                
                dane_handlowca = HANDLOWCY[wybrany_handlowiec]
                
                # WAŻNE: Wstęp na drugiej stronie oferty (wypowiedź właściciela)
                # jest ZAWSZE od Adama Trepki jako CEO It`s Wrap, niezależnie kto
                # operacyjnie przygotowuje ofertę. To strategiczna decyzja - klient
                # dostaje "powitanie od szefa", a wybrany handlowiec dalej widnieje
                # jako kontaktowa osoba do bieżącej obsługi (telefon/email niżej w ofercie).
                AUTOR_WSTEPU_IMIE = "Adam Trepka"
                if AUTOR_WSTEPU_IMIE in HANDLOWCY:
                    AUTOR_WSTEPU_STANOWISKO = HANDLOWCY[AUTOR_WSTEPU_IMIE]["stanowisko"]
                else:
                    # Fallback gdyby Adam Trepka został usunięty/przemianowany w słowniku
                    AUTOR_WSTEPU_STANOWISKO = "CEO It`s Wrap"
                
                # (v2.1) Priorytet ma tekst z pola edycji (jeśli handlowiec go
                # wygenerował/poprawił). Puste pole = automat, jak dotychczas.
                _wstep_z_edytora = st.session_state.get('wstep_editor', '').strip()
                if _wstep_z_edytora:
                    wygenerowany_wstep = _wstep_z_edytora
                else:
                    wygenerowany_wstep = generate_ai_intro_text(
                        klient, final_brand, final_model, pakiet, final_foil_text,
                        AUTOR_WSTEPU_IMIE, AUTOR_WSTEPU_STANOWISKO
                    )

                cena_koncowa_str = f"{cena_koncowa:,.2f} zł".replace(',', 'X').replace('.', ',').replace('X', ' ')
                cena_katalogowa_str = f"{cena_manual:,.2f} zł".replace(',', 'X').replace('.', ',').replace('X', ' ')

                replacements = {
                    "{{KLIENT}}": klient, 
                    "{{MODEL_AUTA}}": f"{final_brand} {final_model}",
                    "{{RODZAJ_FOLII}}": final_foil_text, 
                    "{{USLUGA_NAZWA}}": pakiet,
                    "{{NR_OFERTY}}": nr_o,
                    "{{CENA_KATALOG}}": cena_katalogowa_str,
                    "{{CENA_KONCOWA}}": cena_koncowa_str,
                    "{{WSTEP_AI}}": wygenerowany_wstep,
                    "{{HANDLOWIEC_IMIE}}": wybrany_handlowiec,
                    "{{HANDLOWIEC_TEL}}": dane_handlowca["telefon"],
                    "{{HANDLOWIEC_EMAIL}}": dane_handlowca["email"]
                }

                okladka = next((f for f in pliki_na_dysku if f['name'].startswith('1_')), None)
                wstep_slide = next((f for f in pliki_na_dysku if f['name'].lower().startswith('1b_')), None)
                
                # ==========================================================
                # DOPASOWANIE SZABLONU PRODUKTOWEGO (slajd z opisem folii)
                # ==========================================================
                # Nazwy plików w nowym formacie zaczynają się od daty:
                # - 20251209_PPF_kolor_XPEL      -> XPEL Color (PPF kolorowy)
                # - 20251209_PPF_XPELStealth     -> XPEL Stealth (mat/satyna PPF)
                # - 20251209_PPF_XPELUltimate    -> XPEL Ultimate Plus (połysk PPF)
                # - 20251209_PPF_XPELXtreme      -> XPEL Xtreme (premium PPF)
                # - 20260313_Zmiana koloru_3M    -> 3M wrap
                # - 20260313_Zmiana koloru_AVERY -> Avery wrap
                # - 20260313_Zmiana koloru_ORACLE-> Oracal wrap (w nazwie pliku "Oracle")
                # - 20260430_Zmiana koloru_PWF   -> PWF wrap
                #
                # Wzorzec wyszukiwania: case-insensitive po słowie kluczowym
                # w nazwie pliku, bez polegania na startswith('2_').
                
                def znajdz_szablon(slowo_klucz, dodatkowe_filtry=None):
                    """Szuka pliku zawierającego dane słowo (case-insensitive).
                    dodatkowe_filtry: lista dodatkowych słów, które TEŻ muszą być w nazwie."""
                    for f in pliki_na_dysku:
                        nazwa_lower = f['name'].lower()
                        if slowo_klucz.lower() in nazwa_lower:
                            if dodatkowe_filtry:
                                if all(filt.lower() in nazwa_lower for filt in dodatkowe_filtry):
                                    return f
                            else:
                                return f
                    return None
                
                produkt = None
                
                # Priorytet 1: oferta reklamowa (pakiet o nazwie "reklama")
                if "reklam" in pakiet.lower():
                    produkt = znajdz_szablon('reklama')
                # Priorytet 2: producent folii - mapowanie na słowo w nazwie pliku
                elif f_brand == "3M 2080 Series":
                    produkt = znajdz_szablon('zmiana koloru', dodatkowe_filtry=['3m'])
                elif f_brand == "Avery Dennison SW900":
                    produkt = znajdz_szablon('zmiana koloru', dodatkowe_filtry=['avery'])
                elif f_brand == "Oracal 970RA":
                    # W nazwie pliku jest "ORACLE" (typo z pliku), więc szukamy obu wariantów
                    produkt = znajdz_szablon('zmiana koloru', dodatkowe_filtry=['oracle']) or \
                              znajdz_szablon('zmiana koloru', dodatkowe_filtry=['oracal'])
                elif f_brand == "Arlon":
                    produkt = znajdz_szablon('zmiana koloru', dodatkowe_filtry=['arlon'])
                elif f_brand == "PWF (Platinum Wrapping Film)":
                    # PWF Performance to PPF, reszta linii to wrap kolorowy
                    if "Performance" in f_cat:
                        # Linia Performance łączy PPF z kolorem - jeśli mamy dedykowany szablon, użyj
                        produkt = znajdz_szablon('ppf', dodatkowe_filtry=['pwf']) or \
                                  znajdz_szablon('zmiana koloru', dodatkowe_filtry=['pwf'])
                    else:
                        produkt = znajdz_szablon('zmiana koloru', dodatkowe_filtry=['pwf'])
                # Priorytet 3: XPEL - rozróżnienie po linii produktowej
                elif f_brand == "XPEL (Folie Ochronne PPF)":
                    if "Ultimate" in f_color:
                        produkt = znajdz_szablon('xpelultimate') or znajdz_szablon('ultimate')
                    elif "Stealth" in f_color:
                        produkt = znajdz_szablon('xpelstealth') or znajdz_szablon('stealth')
                    elif "Xtreme" in f_color:
                        produkt = znajdz_szablon('xpelxtreme') or znajdz_szablon('xtreme')
                    elif "Color" in f_cat:
                        # XPEL Color - kolorowy PPF
                        produkt = znajdz_szablon('ppf_kolor') or znajdz_szablon('xpel', dodatkowe_filtry=['kolor'])
                    else:
                        # Ogólny XPEL jeśli nie wpadł w żadną kategorię
                        produkt = znajdz_szablon('xpel')
                
                # Fallback: znajdź dowolny szablon zaczynający się od 2_ (stara konwencja)
                # lub zawierający "PPF" / "Zmiana koloru" w nazwie (nowa konwencja)
                if not produkt:
                    produkt = next((f for f in pliki_na_dysku 
                                    if f['name'].startswith('2_') 
                                    or 'ppf' in f['name'].lower() 
                                    or 'zmiana koloru' in f['name'].lower()), None)
                
                if rabat > 0: 
                    zakres = next((f for f in pliki_na_dysku if f['name'].startswith('3') and 'bezrabatu' not in f['name'].lower()), None)
                else: 
                    zakres = next((f for f in pliki_na_dysku if f['name'].startswith('3') and 'bezrabatu' in f['name'].lower()), None)
                
                if not zakres: 
                    zakres = next((f for f in pliki_na_dysku if f['name'].startswith('3')), None)

                # NOWE: ostatnia strona spersonalizowana pod handlowca
                # Plik z prefixem "6_" i imieniem handlowca w nazwie (np. "6_Oferta_ostatnia_Adam Trepka")
                # Fallback: pierwszy ogólny plik "6_" gdy spersonalizowanej wersji nie ma
                koniec = wybierz_strone_koncowa(pliki_na_dysku, wybrany_handlowiec, list(HANDLOWCY.keys()))

                seq = [okladka, wstep_slide, produkt, zakres] + wybrane_dodatki + [koniec]
                seq = [f for f in seq if f]

                for f_info in seq:
                    prs = Presentation(download_file(service, f_info['id']))
                    for slide in prs.slides:
                        if f_info['name'].startswith('1_'):
                            for shape in list(slide.shapes):
                                if "{{FOTO_AUTA}}" in shape.name or (shape.has_text_frame and "{{FOTO_AUTA}}" in shape.text):
                                    pic = slide.shapes.add_picture(io.BytesIO(st.session_state['ai_img']), shape.left, shape.top, shape.width, shape.height)
                                    slide.shapes._spTree.remove(pic._element)
                                    slide.shapes._spTree.insert(2, pic._element)
                                    shape._element.getparent().remove(shape._element)

                        for shape in slide.shapes:
                            replace_text_in_shape(shape, replacements)

                    # (v2) Unikalne nazwy plików tymczasowych - bezpieczne, gdy
                    # dwóch handlowców generuje oferty w tym samym momencie.
                    tmp_p = f"tmp_{uuid.uuid4().hex[:8]}_{f_info['id']}.pptx"
                    prs.save(tmp_p)
                    pdf = pptx_to_pdf(tmp_p)
                    if pdf: writer.append(pdf); os.remove(tmp_p); os.remove(pdf)

                final_io = io.BytesIO(); writer.write(final_io); final_io.seek(0)
                pdf_bytes = final_io.getvalue()
                nazwa_pliku_wyjsciowego = f"Oferta_{final_brand}_{final_model}_{datetime.now().strftime('%H%M%S')}.pdf"
                
                oauth_drive = get_oauth_drive_service()
                wiz_id = ""
                if oauth_drive:
                    folder_oferty_id = pobierz_lub_stworz_folder_oferty_oauth(oauth_drive)
                    utworzony_link = wgraj_pdf_na_dysk(oauth_drive, folder_oferty_id, nazwa_pliku_wyjsciowego, pdf_bytes)
                    # (v2.4) Zapisujemy też PNG wizualizacji - dzięki temu przy
                    # późniejszej edycji oferty NIE trzeba generować obrazu od nowa.
                    try:
                        _wiz_meta = {
                            'name': f"Wizualizacja_{nr_o.replace('/', '_')}_{datetime.now().strftime('%H%M%S')}.png",
                            'parents': [folder_oferty_id]
                        }
                        _wiz_media = MediaIoBaseUpload(io.BytesIO(st.session_state['ai_img']),
                                                       mimetype='image/png', resumable=True)
                        wiz_id = oauth_drive.files().create(
                            body=_wiz_meta, media_body=_wiz_media, fields='id'
                        ).execute().get('id', "")
                    except Exception:
                        wiz_id = ""
                else:
                    utworzony_link = None
                    st.warning("OAuth Drive niedostępny - PDF nie został zapisany w chmurze. Pobierz plik lokalnie poniżej.")
                
                link_do_zapisu = utworzony_link if utworzony_link else "Błąd uploadu"

                # (v2) Pełne parametry oferty do późniejszej edycji (kolumna DaneOferty)
                dane_oferty_json = json.dumps({
                    "handlowiec": wybrany_handlowiec,
                    "klient": klient,
                    "nr_o": nr_o,
                    "brand": brand,
                    "final_brand": final_brand,
                    "final_model": final_model,
                    "body": body,
                    "segment": segment_final,
                    "year": year,
                    "gen_code": gen_code,
                    "f_brand": f_brand,
                    "f_cat": f_cat,
                    "f_color": f_color,
                    "paint_color": paint_color,
                    "kategoria": kategoria,
                    "pakiet": pakiet,
                    "cena_manual": cena_manual,
                    "rabat": rabat,
                    "dodatki": [d['name'] for d in wybrane_dodatki],
                    # (v2.4) Do odtworzenia przy edycji bez ponownej generacji:
                    "wstep": wygenerowany_wstep,
                    "wizualizacja_id": wiz_id,
                    # (v2.6) Dane klienta z Pipedrive (jeśli osoba istniała):
                    "pd_person_id": (_pd_osoba or {}).get("id", ""),
                    "klient_email": (_pd_osoba or {}).get("email", ""),
                    "klient_telefon": (_pd_osoba or {}).get("telefon", ""),
                    "klient_firma": (_pd_osoba or {}).get("organizacja", ""),
                }, ensure_ascii=False)

                if zapisz_do_rejestru(nr_o, wybrany_handlowiec, klient, f"{final_brand} {final_model}", pakiet, final_foil_text, cena_koncowa, link_do_zapisu, dane_oferty_json):
                    st.success(f"✅ Oferta zapisana w systemie CRM! Plik został zachowany w folderze 'Oferty'.")

                # (v2.3) Zapamiętujemy wygenerowaną ofertę w sesji - dzięki temu
                # przycisk pobierania i sekcja Pipedrive nie znikają po
                # przeładowaniu widoku (Streamlit wykonuje rerun po każdym kliknięciu).
                st.session_state['ostatnia_oferta'] = {
                    "pdf_bytes": pdf_bytes,
                    "nazwa_pliku": nazwa_pliku_wyjsciowego,
                    "nr_o": nr_o,
                    "klient": klient,
                    "auto": f"{final_brand} {final_model}",
                    "pakiet": pakiet,
                    "cena": cena_koncowa,
                    "pipedrive_deal": None,
                    "pd_person_id": (_pd_osoba or {}).get("id", ""),
                }
                st.balloons()

    # ==========================================================
    # (v2.3) SEKCJA PO WYGENEROWANIU OFERTY
    # Pobieranie PDF + dodawanie do Pipedrive NA ŻYCZENIE,
    # z wyborem lejka (pipeline) i etapu (stage).
    # ==========================================================
    _oo = st.session_state.get('ostatnia_oferta')
    if _oo:
        st.markdown("---")
        st.markdown(f"#### 📄 Wygenerowana oferta: {_oo['nr_o']} · {_oo['auto']}"
                    + (f" · {_oo['klient']}" if _oo['klient'].strip() else ""))
        st.download_button("📥 POBIERZ OFERTĘ PDF LOKALNIE",
                           data=_oo['pdf_bytes'],
                           file_name=_oo['nazwa_pliku'],
                           key="dl_ostatnia_oferta")

        st.markdown("##### 📤 Pipedrive - szansa sprzedaży")
        if _oo.get('pipedrive_deal'):
            st.success(f"✅ Ta oferta jest już w Pipedrive - szansa #{_oo['pipedrive_deal']} z załączonym PDF.")
        elif not _pd_dostepny:
            st.caption(
                f"ℹ️ Pipedrive nieaktywny dla: **{wybrany_handlowiec}** - brak tokenu API. "
                f"Token można uzupełnić w zakładce **🛠 Ustawienia** "
                f"(w Pipedrive znajdziesz go w: Ustawienia → Personal preferences → API)."
            )
        else:
            _lejki = pipedrive_pobierz_lejki(_pd_token)
            if not _lejki:
                st.warning("Nie udało się pobrać lejków z Pipedrive - sprawdź, czy token API jest poprawny.")
            else:
                col_pd1, col_pd2, col_pd3 = st.columns([2, 2, 1])
                with col_pd1:
                    _lejki_nazwy = [n for _, n in _lejki]
                    _wybrany_lejek = st.selectbox("Lejek (pipeline)", _lejki_nazwy, key="pd_lejek")
                    _lejek_id = next(i for i, n in _lejki if n == _wybrany_lejek)
                with col_pd2:
                    _etapy = pipedrive_pobierz_etapy(_pd_token, _lejek_id)
                    _stage_id = None
                    if _etapy:
                        _etapy_nazwy = [n for _, n in _etapy]
                        _wybrany_etap = st.selectbox("Etap", _etapy_nazwy, index=0, key=f"pd_etap_{_lejek_id}")
                        _stage_id = next(i for i, n in _etapy if n == _wybrany_etap)
                with col_pd3:
                    st.write("")
                    st.write("")
                    _klik_pd = st.button("📤 DODAJ", key="pd_dodaj")
                if _klik_pd:
                    with st.spinner("Dodaję szansę sprzedaży w Pipedrive..."):
                        # (v2.6) Jeśli osoba została znaleziona już przy wpisywaniu
                        # klienta - wiążemy szansę po jej ID (zero ryzyka duplikatu).
                        if _oo.get('pd_person_id'):
                            _person_id, _osoba_istniala = _oo['pd_person_id'], True
                        elif _oo['klient'].strip():
                            _person_id, _osoba_istniala = pipedrive_znajdz_lub_utworz_osobe(_pd_token, _oo['klient'].strip())
                        else:
                            _person_id, _osoba_istniala = (None, False)
                        _deal_id = pipedrive_utworz_szanse(
                            _pd_token,
                            f"{_oo['nr_o']} | {_oo['auto']} | {_oo['pakiet']}",
                            _person_id,
                            _oo['cena'],
                            stage_id=_stage_id
                        )
                        if _deal_id:
                            if pipedrive_zalacz_pdf(_pd_token, _deal_id, _oo['pdf_bytes'], _oo['nazwa_pliku']):
                                _oo['pipedrive_deal'] = _deal_id
                                st.session_state['ostatnia_oferta'] = _oo
                                st.success(
                                    f"📤 Szansa #{_deal_id} utworzona w lejku '{_wybrany_lejek}' "
                                    f"({'istniejąca osoba' if _osoba_istniala else 'nowa osoba'}), PDF załączony."
                                )
                            else:
                                st.warning(f"Szansa #{_deal_id} utworzona, ale załączenie PDF nie powiodło się.")
                        else:
                            st.warning("Nie udało się utworzyć szansy sprzedaży - sprawdź token API.")

with tab_rejestr:
    st.markdown("""
    <div class="iw-main-header">
        <div class="iw-main-header-text">
            <div class="iw-main-header-title">Ewidencja Ofert</div>
            <div class="iw-main-header-subtitle">Rejestr · Archiwum · Kontrola</div>
        </div>
    </div>
    """, unsafe_allow_html=True)
    try:
        dane_rejestru = pobierz_rejestr()
        if not dane_rejestru:
            st.info("Rejestr jest pusty lub arkusz nie zawiera jeszcze wpisów.")
        else:
            # (v2.4) JEDNA lista ofert. Handlowiec widzi swoje, admin wszystkie
            # (role działają dopiero, gdy istnieje baza handlowców z rolami).
            if jest_adminem:
                rekordy = list(dane_rejestru)
            else:
                rekordy = [r for r in dane_rejestru
                           if _handlowiec_wiersza(r) == wybrany_handlowiec]

            if not rekordy:
                st.info("Brak ofert przypisanych do tego handlowca.")
            else:
                col_f1, col_f2 = st.columns([3, 1])
                with col_f1:
                    szukaj = st.text_input("🔎 Szukaj (klient, numer oferty, auto...)", key="rej_szukaj")
                with col_f2:
                    st.write("")
                    # CSV bez technicznej kolumny z JSON-em
                    _df_csv = pd.DataFrame(rekordy)
                    _df_csv = _df_csv[[c for c in _df_csv.columns
                                       if re.sub(r'\W', '', str(c)).lower() != 'daneoferty']]
                    st.download_button(
                        label="⬇️ CSV",
                        data=_df_csv.to_csv(index=False).encode('utf-8'),
                        file_name=f"Rejestr_Ofert_ITSWRAP_{datetime.now().strftime('%Y%m%d')}.csv",
                        mime='text/csv',
                        use_container_width=True,
                    )

                if szukaj.strip():
                    _s = szukaj.strip().lower()
                    rekordy = [r for r in rekordy
                               if _s in ' '.join(str(v) for v in r.values()).lower()]

                LIMIT_LISTY = 50
                widoczne = list(reversed(rekordy))[:LIMIT_LISTY]
                st.caption(f"Ofert: {len(rekordy)} · najnowsze u góry"
                           + (f" · pokazuję {LIMIT_LISTY} - zawęź wyszukiwarką" if len(rekordy) > LIMIT_LISTY else ""))

                for _i, r in enumerate(widoczne):
                    _json_oferty = _wyciagnij_dane_oferty(r)
                    _link_pdf = next((str(v).strip() for v in r.values()
                                      if str(v).strip().startswith('http')), "")
                    c1, c2, c3, c4 = st.columns([2.4, 3.0, 1.0, 1.0])
                    with c1:
                        st.markdown(
                            f"**{_pole(r, 'Nr oferty', 'Numer oferty', 'Nr') or '(bez numeru)'}**  \n"
                            f"<span style='color:#4A5568; font-size:0.85rem'>"
                            f"{_pole(r, 'Data')} · {_handlowiec_wiersza(r)}</span>",
                            unsafe_allow_html=True)
                    with c2:
                        st.markdown(
                            f"{_pole(r, 'Klient') or '(klient bez nazwy)'}  \n"
                            f"<span style='color:#4A5568; font-size:0.85rem'>"
                            f"{_pole(r, 'Auto', 'Samochód')} · {_pole(r, 'Usługa', 'Usluga')} · "
                            f"{_pole(r, 'Cena', 'Cena netto', 'Cena końcowa')}</span>",
                            unsafe_allow_html=True)
                    with c3:
                        if _link_pdf:
                            st.link_button("📥 Pobierz", _link_pdf, use_container_width=True)
                        else:
                            st.caption("brak PDF")
                    with c4:
                        if _json_oferty:
                            if st.button("✏️ Edytuj", key=f"rej_edit_{_i}", use_container_width=True):
                                try:
                                    _d = json.loads(_json_oferty)
                                    st.session_state['edycja'] = _d
                                    # Wstęp: wczyta się do edytora przy następnym przebiegu
                                    st.session_state['zaladuj_wstep'] = _d.get('wstep', '')
                                    # Wizualizacja: jeśli była zapisana na Dysku - pobieramy,
                                    # żeby NIE trzeba było generować jej od nowa.
                                    st.session_state.pop('ai_img', None)
                                    _wiz_id = _d.get('wizualizacja_id', '')
                                    if _wiz_id:
                                        with st.spinner("Wczytuję zapisaną wizualizację..."):
                                            try:
                                                _oauth_tmp = get_oauth_drive_service()
                                                if _oauth_tmp:
                                                    st.session_state['ai_img'] = download_file(_oauth_tmp, _wiz_id).getvalue()
                                            except Exception:
                                                pass
                                    st.rerun()
                                except Exception as e:
                                    st.error(f"Nie udało się wczytać oferty: {e}")
                        else:
                            st.caption("✏️ n/d")
                    st.markdown("<hr style='margin:2px 0 10px 0; opacity:0.2'>", unsafe_allow_html=True)
    except Exception as e:
        st.warning(f"Brak możliwości wczytania rejestru. Błąd: {e}")

with tab_ustawienia:
    st.markdown("""
    <div class="iw-main-header">
        <div class="iw-main-header-text">
            <div class="iw-main-header-title">Ustawienia</div>
            <div class="iw-main-header-subtitle">Folie · Handlowcy · Konfiguracja</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # ==========================================================
    # (v2.3) FOLIE - dodawanie kolorów z poziomu aplikacji
    # ==========================================================
    st.markdown("### 🎨 Baza folii")

    if not FOLIE_Z_ARKUSZA:
        st.info(
            "Baza kolorów działa obecnie na liście wbudowanej w kod. Kliknij poniżej, aby "
            "przenieść ją do arkusza Google - od tego momentu nowe kolory dodaje się "
            "formularzem poniżej, bez dotykania kodu i bez wchodzenia do arkusza."
        )
        if st.button("📥 UTWÓRZ BAZĘ FOLII I ZAIMPORTUJ OBECNE KOLORY", disabled=not zalogowany):
            with st.spinner("Tworzę zakładkę 'Folie' i importuję bazę kolorów..."):
                try:
                    utworz_zakladke_folie_i_importuj()
                    st.success("✅ Baza folii przeniesiona do arkusza. Odświeżam...")
                    st.rerun()
                except Exception as e:
                    st.error(f"Nie udało się utworzyć bazy: {e}")
    else:
        st.caption(f"✅ Baza kolorów działa z arkusza. Producentów: "
                   f"{len([p for p in FOIL_GROUPS if p != 'Inne (wpisz ręcznie)'])}, "
                   f"kolorów: {sum(len(k) for p, c in FOIL_GROUPS.items() if p != 'Inne (wpisz ręcznie)' for k in c.values())}.")

        st.markdown("#### ➕ Dodaj nowy kolor")
        _prod_opcje = [p for p in FOIL_GROUPS.keys() if p != "Inne (wpisz ręcznie)"] + ["➕ Nowy producent..."]
        _sel_prod = st.selectbox("Producent", _prod_opcje, key="ust_prod")

        _kategorie_nowego = ""
        if _sel_prod == "➕ Nowy producent...":
            _nowy_prod = st.text_input("Nazwa nowego producenta", key="ust_nowy_prod",
                                       placeholder="np. Hexis Skintac")
            _kat_wybor = st.multiselect(
                "Kategorie cennika dla tego producenta (puste = wszystkie)",
                ["PPF", "Zmiana koloru"], key="ust_kat_nowego")
            _kategorie_nowego = ", ".join(_kat_wybor)
            _final_prod = _nowy_prod.strip()
            _wyk_opcje = ["➕ Nowe wykończenie..."]
        else:
            _final_prod = _sel_prod
            _wyk_opcje = list(FOIL_GROUPS[_sel_prod].keys()) + ["➕ Nowe wykończenie..."]

        _sel_wyk = st.selectbox("Wykończenie", _wyk_opcje, key="ust_wyk")
        if _sel_wyk == "➕ Nowe wykończenie...":
            _final_wyk = st.text_input("Nazwa nowego wykończenia", key="ust_nowe_wyk",
                                       placeholder="np. Połysk (Gloss)").strip()
        else:
            _final_wyk = _sel_wyk

        _nowy_kolor = st.text_input(
            "Nazwa koloru",
            key="ust_kolor",
            placeholder="np. Zielony Jasny Połysk (G16) - Gloss Light Green",
            help="Zalecany format: Nazwa Polska (Kod) - Nazwa Angielska. "
                 "Nazwa angielska jest używana w promptach do AI."
        )

        if st.button("💾 DODAJ KOLOR DO BAZY", disabled=not zalogowany, key="ust_dodaj_kolor"):
            if not _final_prod or not _final_wyk or not _nowy_kolor.strip():
                st.error("Uzupełnij producenta, wykończenie i nazwę koloru.")
            else:
                _duplikat = _nowy_kolor.strip() in FOIL_GROUPS.get(_final_prod, {}).get(_final_wyk, [])
                if _duplikat:
                    st.warning("Ten kolor już istnieje w bazie dla tego producenta i wykończenia.")
                else:
                    try:
                        dodaj_kolor_folii(_final_prod, _final_wyk, _nowy_kolor, _kategorie_nowego)
                        st.success(f"✅ Dodano: {_final_prod} → {_final_wyk} → {_nowy_kolor.strip()}")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Nie udało się dodać koloru: {e}")

    st.markdown("---")

    # ==========================================================
    # (v2.3) HANDLOWCY - zarządzanie z poziomu aplikacji (admin)
    # ==========================================================
    st.markdown("### 👥 Handlowcy")
    if not jest_adminem:
        st.caption("Zarządzanie handlowcami jest dostępne tylko dla administratora.")
    else:
        if not HANDLOWCY_Z_ARKUSZA:
            st.info(
                "Lista handlowców działa obecnie na danych wbudowanych w kod - bez PIN-ów, "
                "tokenów Pipedrive i ról. Kliknij poniżej, aby przenieść ją do arkusza - od tego "
                "momentu wszystko (PIN-y, tokeny, role, nowi handlowcy) ustawiasz formularzem poniżej."
            )
            if st.button("📥 UTWÓRZ BAZĘ HANDLOWCÓW I ZAIMPORTUJ OBECNĄ LISTĘ", disabled=not zalogowany):
                with st.spinner("Tworzę zakładkę 'Handlowcy' i importuję listę..."):
                    try:
                        utworz_zakladke_handlowcy_i_importuj(HANDLOWCY)
                        st.success("✅ Lista handlowców przeniesiona do arkusza. Odświeżam...")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Nie udało się utworzyć bazy: {e}")
        else:
            _df_h = pd.DataFrame([
                {
                    "Handlowiec": nazwa,
                    "Stanowisko": d.get("stanowisko", ""),
                    "Telefon": d.get("telefon", ""),
                    "Email": d.get("email", ""),
                    "PIN": "•••" if d.get("pin") else "(brak)",
                    "Pipedrive": "✅" if d.get("pipedrive_token") else "❌",
                    "Rola": d.get("rola", "handlowiec"),
                }
                for nazwa, d in HANDLOWCY.items()
            ])
            st.dataframe(_df_h, hide_index=True, use_container_width=True)

            st.markdown("#### ✏️ Dodaj lub edytuj handlowca")
            _h_opcje = list(HANDLOWCY.keys()) + ["➕ Nowy handlowiec..."]
            _sel_h = st.selectbox("Handlowiec", _h_opcje, key="ust_h_sel")

            if _sel_h == "➕ Nowy handlowiec...":
                _h_dane = {"stanowisko": "", "telefon": "", "email": "",
                           "pin": "", "pipedrive_token": "", "rola": "handlowiec"}
                _h_nazwa = st.text_input("Imię i nazwisko", key="ust_h_nazwa")
            else:
                _h_dane = HANDLOWCY[_sel_h]
                _h_nazwa = _sel_h

            col_h1, col_h2 = st.columns(2)
            with col_h1:
                _h_stanowisko = st.text_input("Stanowisko", value=_h_dane.get("stanowisko", ""), key=f"ust_h_stan_{_sel_h}")
                _h_telefon = st.text_input("Telefon", value=_h_dane.get("telefon", ""), key=f"ust_h_tel_{_sel_h}")
                _h_email = st.text_input("Email", value=_h_dane.get("email", ""), key=f"ust_h_mail_{_sel_h}")
            with col_h2:
                _h_pin = st.text_input("PIN (puste = bez logowania)", value=_h_dane.get("pin", ""), key=f"ust_h_pin_{_sel_h}")
                _h_token = st.text_input("Token API Pipedrive", value=_h_dane.get("pipedrive_token", ""),
                                         type="password", key=f"ust_h_tok_{_sel_h}",
                                         help="Pipedrive: Ustawienia → Personal preferences → API")
                _h_rola = st.selectbox("Rola", ["handlowiec", "admin"],
                                       index=(1 if _h_dane.get("rola") == "admin" else 0),
                                       key=f"ust_h_rola_{_sel_h}")

            if st.button("💾 ZAPISZ HANDLOWCA", key="ust_h_zapisz"):
                if not _h_nazwa.strip():
                    st.error("Podaj imię i nazwisko handlowca.")
                else:
                    try:
                        _wynik = zapisz_handlowca(_h_nazwa, _h_stanowisko, _h_telefon,
                                                  _h_email, _h_pin, _h_token, _h_rola)
                        st.success(f"✅ {'Zaktualizowano dane' if _wynik == 'zaktualizowano' else 'Dodano handlowca'}: {_h_nazwa.strip()}")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Nie udało się zapisać: {e}")

st.markdown("""
<div class="iw-footer">
    <strong>IT'S WRAP</strong> · ONDRE.PL · Rynek Śródecki 4, 61-126 Poznań · +48 602 494 133
    <br/>
    <span style="opacity: 0.6; letter-spacing: 0.03em;">www.itswrap.pl · Professional Car Wrapping &amp; PPF · Certyfikowana jakość 3M</span>
</div>
""", unsafe_allow_html=True)
